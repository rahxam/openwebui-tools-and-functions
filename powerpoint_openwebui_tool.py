"""
title: PowerPoint Creator
author: Maximilian Hartig
author_url: https://github.com/rahxam/aequitas-confluence-calendar
description: Create and manage PowerPoint presentations with slides, images, tables, and SVG graphics
required_open_webui_version: 0.4.0
requirements: python-pptx>=0.6.21
version: 1.0.0
licence: MIT
"""

import os
import re
import uuid
import base64
import io
from typing import List,  Optional
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN


class Tools:
    def __init__(self):
        """Initialize the PowerPoint Tool."""
        self.presentations = {}

    def _generate_presentation_id(self, title: str) -> str:
        cleaned_title = re.sub(r"[^\w\s\-]", "", title)
        cleaned_title = re.sub(r"[\s\-]+", "_", cleaned_title)
        cleaned_title = cleaned_title.strip("_").lower()
        if len(cleaned_title) > 20:
            cleaned_title = cleaned_title[:20].rstrip("_")
        short_uuid = uuid.uuid4().hex[:8]

        if cleaned_title:
            return f"{cleaned_title}_{short_uuid}"
        else:
            return f"presentation_{short_uuid}"

    async def create_presentation(self, title: str, __event_emitter__=None) -> str:
        """
        Create a new empty PowerPoint presentation with a title. Do not forget to create the individual slides afterwards. Once you are done, call the save tool.

        :param title: The title/name for the presentation
        """
        try:
            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Creating presentation: {title}",
                            "done": False,
                        },
                    }
                )

            template_path = os.path.join('.', 'data', 'templates', 'aequitas_master.pptx')
            if os.path.exists(template_path):
                prs = Presentation(template_path)
            else:
                prs = Presentation()
            prs_id = self._generate_presentation_id(title)

            # Ensure uniqueness (though UUID makes collision extremely unlikely)
            while prs_id in self.presentations:
                prs_id = self._generate_presentation_id(title)

            self.presentations[prs_id] = prs

            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Presentation '{prs_id}' created successfully",
                            "done": True,
                        },
                    }
                )

            return f"Created presentation: {prs_id}"

        except Exception as e:
            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Error creating presentation: {str(e)}",
                            "done": True,
                        },
                    }
                )
            return f"Error creating presentation: {str(e)}"

    async def add_title_slide(
        self,
        prs_id: str,
        title: str,
        subtitle: Optional[str] = None,
        __event_emitter__=None,
    ) -> str:
        """
        Add a title slide to the presentation.

        :param prs_id: The presentation ID
        :param title: The title text
        :param subtitle: Optional subtitle text
        """
        try:
            # Wenn nur eine Präsentation geladen ist, nimm diese, egal was übergeben wurde
            if len(self.presentations) == 1:
                prs_id = list(self.presentations.keys())[0]
            if prs_id not in self.presentations:
                return f"Error: Presentation '{prs_id}' not found"

            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Adding title slide to {prs_id}",
                            "done": False,
                        },
                    }
                )

            prs = self.presentations[prs_id]
            slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title_shape = slide.placeholders[10]
            if title_shape is not None:
                title_shape.text = title

            # Set subtitle if provided
            if (
                subtitle
                and hasattr(slide, "placeholders")
                and len(slide.placeholders) > 1
            ):
                subtitle_shape = slide.placeholders[14]
                if subtitle_shape is not None and hasattr(subtitle_shape, "has_text_frame") and subtitle_shape.has_text_frame:
                    subtitle_shape.text = subtitle

            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Title slide added successfully",
                            "done": True,
                        },
                    }
                )

            return f"Added title slide at position {len(prs.slides)}"

        except Exception as e:
            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Error adding title slide: {str(e)}",
                            "done": True,
                        },
                    }
                )
            return f"Error adding title slide: {str(e)}"

    async def add_content_slide(
        self, prs_id: str, title: str,  subtitle: str, content: List[str], __event_emitter__=None
    ) -> str:
        """
        Add a content slide with bullet points.

        :param prs_id: The presentation ID
        :param title: The slide title
        :param subtitle: The slide subtitle (can be empty)
        :param content: List of bullet points
        """
        try:
            if len(self.presentations) == 1:
                prs_id = list(self.presentations.keys())[0]
            if prs_id not in self.presentations:
                return f"Error: Presentation '{prs_id}' not found"

            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Adding content slide: {title}",
                            "done": False,
                        },
                    }
                )

            prs = self.presentations[prs_id]
            slide_layout = prs.slide_layouts[8]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title_shape = slide.placeholders[11]
            if title_shape is not None:
                title_shape.text = title

            # Set subtitle if provided
            if (
                subtitle
                and hasattr(slide, "placeholders")
                and len(slide.placeholders) > 1
            ):
                subtitle_shape = slide.placeholders[15]
                if subtitle_shape is not None and hasattr(subtitle_shape, "has_text_frame") and subtitle_shape.has_text_frame:
                    subtitle_shape.text = subtitle
                

            # Find content placeholder
            content_placeholder = slide.placeholders[14]

            # If no content placeholder found, add a textbox
            if not content_placeholder:
                content_placeholder = slide.shapes.add_textbox(
                    Inches(1), Inches(2), Inches(8), Inches(4)
                )

            # Add bullet points
            tf = content_placeholder.text_frame
            tf.clear()  # Clear any existing text

            # Add first bullet point
            if content and len(content) > 0:
                p = tf.paragraphs[0]
                p.text = content[0]
                p.level = 0

                # Add remaining bullet points
                for bullet_text in content[1:]:
                    p = tf.add_paragraph()
                    p.text = bullet_text
                    p.level = 0

            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Content slide added with {len(content)} bullet points",
                            "done": True,
                        },
                    }
                )

            return f"Added content slide at position {len(prs.slides)}"

        except Exception as e:
            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Error adding content slide: {str(e)}",
                            "done": True,
                        },
                    }
                )
            return f"Error adding content slide: {str(e)}"

    async def add_section_slide(
        self,
        prs_id: str,
        section_title: str,
        section_subtitle: str,
        __event_emitter__=None,
    ) -> str:
        """
        Add a section divider slide with a large title.

        :param prs_id: The presentation ID
        :param section_title: The section title text
        :param section_subtitle: The section subtitle text (can be empty)
        """
        try:
            if len(self.presentations) == 1:
                prs_id = list(self.presentations.keys())[0]
            if prs_id not in self.presentations:
                return f"Error: Presentation '{prs_id}' not found"

            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Adding section slide: {section_title}",
                            "done": False,
                        },
                    }
                )

            prs = self.presentations[prs_id]
            slide_layout = prs.slide_layouts[1]  # Blank slide layout
            slide = prs.slides.add_slide(slide_layout)
           
            # Set title
            title_shape = slide.placeholders[10]
            if title_shape is not None:
                title_shape.text = section_title

            subtitle_shape = slide.placeholders[14]
            if subtitle_shape is not None:
                subtitle_shape.text = section_subtitle

            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Section slide added successfully",
                            "done": True,
                        },
                    }
                )

            return f"Added section slide at position {len(prs.slides)}"

        except Exception as e:
            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Error adding section slide: {str(e)}",
                            "done": True,
                        },
                    }
                )
            return f"Error adding section slide: {str(e)}"

    async def add_table_slide(
        self,
        prs_id: str,
        title: str,
        subtitle: str,
        headers: List[str],
        rows: List[List[str]],
        __event_emitter__=None,
    ) -> str:
        """
        Add a slide with a table.

        :param prs_id: The presentation ID
        :param title: The slide title
        :param subtitle: The slide subtitle (can be empty)
        :param headers: List of column headers
        :param rows: List of rows, where each row is a list of cell values
        """
        try:
            if len(self.presentations) == 1:
                prs_id = list(self.presentations.keys())[0]
            if prs_id not in self.presentations:
                return f"Error: Presentation '{prs_id}' not found"

            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Adding table slide: {title}",
                            "done": False,
                        },
                    }
                )

            prs = self.presentations[prs_id]
            slide_layout = prs.slide_layouts[14]  # Blank slide layout
            slide = prs.slides.add_slide(slide_layout)

            
           
            # Set title
            title_shape = slide.placeholders[11]
            if title_shape is not None:
                title_shape.text = title
                
            subtitle_shape = slide.placeholders[15]
            if subtitle_shape is not None:
                subtitle_shape.text = subtitle

            # Calculate table dimensions
            num_rows = len(rows) + 1  # +1 for header row
            num_cols = len(headers)

            # Create table
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(12)
            height = Inches(0.5 * num_rows)

            table = slide.shapes.add_table(
                num_rows, num_cols, left, top, width, height
            ).table

            # Set column widths
            col_width = int(width / num_cols)
            for i in range(num_cols):
                table.columns[i].width = col_width

            # Add headers
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Add data rows
            for i, row in enumerate(rows):
                for j, cell_value in enumerate(row):
                    if j < num_cols:
                        cell = table.cell(i + 1, j)
                        cell.text = str(cell_value)

            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Table slide added with {len(rows)} rows",
                            "done": True,
                        },
                    }
                )

            return f"Added table slide at position {len(prs.slides)}"

        except Exception as e:
            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Error adding table slide: {str(e)}",
                            "done": True,
                        },
                    }
                )
            return f"Error adding table slide: {str(e)}"

    def get_presentation_info(self, prs_id: str, __event_emitter__=None) -> str:
        """
        Get information about the presentation.

        :param prs_id: The presentation ID
        """
        try:
            if len(self.presentations) == 1:
                prs_id = list(self.presentations.keys())[0]
            if prs_id not in self.presentations:
                return f"Error: Presentation '{prs_id}' not found"

            prs = self.presentations[prs_id]
            slide_count = len(prs.slides)

            info = f"Presentation: {prs_id}\n"
            info += f"Number of slides: {slide_count}\n"
            info += "Slide layouts available:\n"

            for i, layout in enumerate(prs.slide_layouts):
                layout_name = layout.name if hasattr(layout, "name") else f"Layout {i}"
                info += f"  - {layout_name}\n"

            return info

        except Exception as e:
            return f"Error getting presentation info: {str(e)}"

    async def save_presentation_base64(
        self, prs_id: str, __event_emitter__=None
    ) -> str:
        """
        Save the presentation as pptx and return it as base64 for download.

        :param prs_id: The presentation ID
        """
        try:
            if len(self.presentations) == 1:
                prs_id = list(self.presentations.keys())[0]
            elif prs_id is None:
                prs_id = next(iter(self.presentations), None)
            if prs_id not in self.presentations:
                return f"Error: Presentation '{prs_id}' not found"

            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Preparing presentation for download",
                            "done": False,
                        },
                    }
                )

            prs = self.presentations[prs_id]

            # Save to a bytes buffer
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)

            # Encode to base64
            base64_data = base64.b64encode(buffer.read()).decode("utf-8")
            file_size = len(buffer.getvalue()) / 1024

            download_uri = f"data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{base64_data}"

            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "message",
                        "data": {
                            "content": f"""
## PowerPoint Presentation Ready for Download

**Presentation ID:** {prs_id}  
**File Size:** {file_size:.2f} KB  
**Slides:** {len(prs.slides)}

### Download Link:
[📥 Download {prs_id}.pptx]({download_uri})

"""
                        },
                    }
                )

            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Presentation ready for download",
                            "done": True,
                        },
                    }
                )

            return f"Presentation '{prs_id}' prepared for download. File size: {file_size:.2f} KB"

        except Exception as e:
            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Error preparing download: {str(e)}",
                            "done": True,
                        },
                    }
                )
            return f"Error preparing presentation for download: {str(e)}"

    async def create_quick_presentation(
        self,
        title: str,
        slide_titles: List[str],
        content_per_slide: List[List[str]],
        __event_emitter__=None,
    ) -> str:
        """
        Quickly create a standard content presentation with multiple slides.

        :param title: Presentation title
        :param slide_titles: List of slide titles
        :param content_per_slide: List of bullet points for each slide
        """
        try:
            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Creating quick presentation: {title}",
                            "done": False,
                        },
                    }
                )

            # Create presentation from template if available
            template_path = os.path.join('.', 'data', 'templates', 'aequitas_master.pptx')
            if os.path.exists(template_path):
                prs = Presentation(template_path)
            else:
                prs = Presentation()
            prs_id = self._generate_presentation_id(title)
            self.presentations[prs_id] = prs

            # Title slide
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title if hasattr(slide.shapes, "title") else None
            if title_shape is not None:
                title_shape.text = title

            # Add subtitle if available
            if hasattr(slide, "placeholders") and len(slide.placeholders) > 1:
                subtitle_shape = slide.placeholders[1]
                if subtitle_shape is not None and hasattr(subtitle_shape, "has_text_frame") and subtitle_shape.has_text_frame:
                    subtitle_shape.text = "Generated Presentation"

            # Content slides
            for i, slide_title in enumerate(slide_titles):
                slide_layout = prs.slide_layouts[1]  # Title and content layout
                slide = prs.slides.add_slide(slide_layout)

                # Set title
                title_shape = slide.shapes.title if hasattr(slide.shapes, "title") else None
                if title_shape is not None:
                    title_shape.text = slide_title

                # Add content if available
                content = content_per_slide[i] if i < len(content_per_slide) else []
                if content:
                    # Find content placeholder
                    content_placeholder = None
                    for shape in slide.placeholders:
                        if shape.placeholder_format.type == 2:  # Body placeholder
                            content_placeholder = shape
                            break

                    # If no content placeholder found, add a textbox
                    if not content_placeholder:
                        content_placeholder = slide.shapes.add_textbox(
                            Inches(1), Inches(2), Inches(8), Inches(4)
                        )

                    if content_placeholder:
                        tf = content_placeholder.text_frame
                        tf.clear()

                        # Add bullet points
                        if content and len(content) > 0:
                            p = tf.paragraphs[0]
                            p.text = content[0]
                            p.level = 0

                            for bullet_text in content[1:]:
                                p = tf.add_paragraph()
                                p.text = bullet_text
                                p.level = 0

            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Quick presentation created with {len(slide_titles)} slides",
                            "done": True,
                        },
                    }
                )

            return f"Created presentation '{prs_id}' with {len(slide_titles)} content slides"

        except Exception as e:
            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Error creating quick presentation: {str(e)}",
                            "done": True,
                        },
                    }
                )
            return f"Error creating quick presentation: {str(e)}"

    def list_presentations(self, __event_emitter__=None) -> str:
        """
        List all active presentations in memory.
        """
        try:
            if not self.presentations:
                return "No presentations currently in memory."

            info = "Active Presentations:\n"
            for prs_id, prs in self.presentations.items():
                slide_count = len(prs.slides)
                info += f"- {prs_id}: {slide_count} slides\n"

            return info

        except Exception as e:
            return f"Error listing presentations: {str(e)}"

    async def remove_presentation(self, prs_id: str, __event_emitter__=None) -> str:
        """
        Remove a presentation from memory.

        :param prs_id: The presentation ID to remove
        """
        try:
            if len(self.presentations) == 1:
                prs_id = list(self.presentations.keys())[0]
            if prs_id not in self.presentations:
                return f"Error: Presentation '{prs_id}' not found"

            del self.presentations[prs_id]

            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Presentation '{prs_id}' removed from memory",
                            "done": True,
                        },
                    }
                )

            return f"Presentation '{prs_id}' removed from memory"

        except Exception as e:
            return f"Error removing presentation: {str(e)}"
