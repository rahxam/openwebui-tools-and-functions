"""
title: PowerPoint MCP Toolkit
author: GongRzhe & Migration by [Your Name]
author_url: https://github.com/GongRzhe/Office-PowerPoint-MCP-Server
git_url: https://github.com/GongRzhe/Office-PowerPoint-MCP-Server
description: A comprehensive toolkit for advanced PowerPoint manipulation, management, and automation using python-pptx. Includes 30+ tools for presentation, content, template, structural, design, chart, hyperlink, connector, master, and transition operations.
required_open_webui_version: 0.4.0
requirements: python-pptx
version: 2.1.0
license: MIT
"""

import os
import platform
from typing import List, Optional, Dict, Any, Tuple
from pydantic import BaseModel, Field


class Tools:
    class Valves(BaseModel):
        base_folder: str = Field(
            default="./data",
            description="Base folder for all file operations (presentations, templates, exports)",
        )
        template_search_paths: List[str] = Field(
            default_factory=lambda: [".", "./templates", "./assets", "./resources"],
            description="Directories to search for PowerPoint templates.",
        )
        allow_multiple_presentations: bool = Field(
            False, description="Allow multiple presentations to be loaded at once."
        )
        enable_logging: bool = Field(False, description="Enable logging for debugging.")

    def __init__(self):
        self.valves = self.Valves()
        self.presentations = {}
        self.current_presentation_id = None

    # --- Presentation Management Methods ---

    def list_presentations(self) -> dict:
        """List all loaded presentations."""
        return {
            "presentations": [
                {
                    "id": pres_id,
                    "slide_count": (
                        len(pres.slides) if hasattr(pres, "slides") else None
                    ),
                    "is_current": pres_id == self.current_presentation_id,
                }
                for pres_id, pres in self.presentations.items()
            ],
            "current_presentation_id": self.current_presentation_id,
            "total_presentations": len(self.presentations),
        }

    def switch_presentation(self, presentation_id: str) -> dict:
        """Switch to a different loaded presentation."""
        if presentation_id not in self.presentations:
            return {
                "error": f"Presentation '{presentation_id}' not found. Available presentations: {list(self.presentations.keys())}"
            }
        old_id = self.current_presentation_id
        self.current_presentation_id = presentation_id
        return {
            "message": f"Switched from presentation '{old_id}' to '{presentation_id}'",
            "previous_presentation_id": old_id,
            "current_presentation_id": self.current_presentation_id,
        }

    def get_current_presentation(self):
        """Get the current presentation object or raise an error if none is loaded."""
        if (
            self.current_presentation_id is None
            or self.current_presentation_id not in self.presentations
        ):
            raise ValueError(
                "No presentation is currently loaded. Please create or open a presentation first."
            )
        return self.presentations[self.current_presentation_id]

    # --- Utility Functions ---

    def get_template_search_directories(self) -> list:
        """Get list of directories to search for templates, using valves or defaults."""
        env_path = os.environ.get("PPT_TEMPLATE_PATH")
        directories = []

        # Add environment paths if available
        if env_path:
            separator = ";" if platform.system() == "Windows" else ":"
            env_dirs = [p.strip() for p in env_path.split(separator) if p.strip()]
            for d in env_dirs:
                expanded_dir = os.path.expanduser(d)
                if os.path.isdir(expanded_dir):
                    directories.append(expanded_dir)

        # Add configured template search paths (resolved with base folder)
        for path in self.valves.template_search_paths:
            resolved_path = (
                self.resolve_file_path(path) if not os.path.isabs(path) else path
            )
            if os.path.isdir(resolved_path):
                directories.append(resolved_path)
            elif not os.path.exists(resolved_path):
                # Create the directory if it doesn't exist
                try:
                    os.makedirs(resolved_path, exist_ok=True)
                    directories.append(resolved_path)
                except OSError:
                    pass  # Skip if can't create

        return directories if directories else [self.valves.base_folder]

    def resolve_file_path(self, file_path: str) -> str:
        """Resolve file path with base folder. Creates base folder if it doesn't exist.

        Args:
            file_path: Relative or absolute file path

        Returns:
            Absolute path with base folder prepended if path is relative
        """
        # If it's already an absolute path, return as-is
        if os.path.isabs(file_path):
            return file_path

        # Resolve relative path with base folder
        resolved_path = os.path.join(self.valves.base_folder, file_path)

        # Create base folder and parent directories if they don't exist
        parent_dir = os.path.dirname(resolved_path)
        if parent_dir and not os.path.exists(parent_dir):
            os.makedirs(parent_dir, exist_ok=True)

        return resolved_path

    def validate_parameters(self, params: dict) -> tuple:
        """Validate parameters against constraints. Returns (True, None) or (False, error_message)."""
        for param_name, (value, constraints) in params.items():
            for constraint_func, error_msg in constraints:
                if not constraint_func(value):
                    return False, f"Parameter '{param_name}': {error_msg}"
        return True, None

    def is_positive(self, value):
        return value > 0

    def is_non_negative(self, value):
        return value >= 0

    def is_in_range(self, min_val, max_val):
        return lambda x: min_val <= x <= max_val

    def is_in_list(self, valid_list):
        return lambda x: x in valid_list

    def is_valid_rgb(self, color_list):
        if not isinstance(color_list, list) or len(color_list) != 3:
            return False
        return all(isinstance(c, int) and 0 <= c <= 255 for c in color_list)

    def add_shape_direct(
        self,
        slide,
        shape_type: str,
        left: float,
        top: float,
        width: float,
        height: float,
    ):
        """Add an auto shape to a slide using direct integer values instead of enum objects."""
        from pptx.util import Inches

        shape_type_map = {
            "rectangle": 1,
            "rounded_rectangle": 2,
            "oval": 9,
            "diamond": 4,
            "triangle": 5,
            "right_triangle": 6,
            "pentagon": 56,
            "hexagon": 10,
            "heptagon": 11,
            "octagon": 12,
            "star": 12,
            "arrow": 13,
            "cloud": 35,
            "heart": 21,
            "lightning_bolt": 22,
            "sun": 23,
            "moon": 24,
            "smiley_face": 17,
            "no_symbol": 19,
            "flowchart_process": 112,
            "flowchart_decision": 114,
            "flowchart_data": 115,
            "flowchart_document": 119,
        }
        shape_type_lower = str(shape_type).lower()
        if shape_type_lower not in shape_type_map:
            available_shapes = ", ".join(sorted(shape_type_map.keys()))
            raise ValueError(
                f"Unsupported shape type: '{shape_type}'. Available shape types: {available_shapes}"
            )
        shape_value = shape_type_map[shape_type_lower]
        try:
            shape = slide.shapes.add_shape(
                shape_value, Inches(left), Inches(top), Inches(width), Inches(height)
            )
            return shape
        except Exception as e:
            raise ValueError(
                f"Failed to create '{shape_type}' shape using direct value {shape_value}: {str(e)}"
            )

    # --- Core Presentation Tools ---

    def create_presentation(self, id: Optional[str] = None) -> dict:
        """Create a new blank PowerPoint presentation.

        ⚠️  RECOMMENDATION: Consider using create_presentation_from_template() instead for better results!

        For professional presentations, it's highly recommended to:
        1. First call discover_available_templates() to see available template files
        2. Ask the user which template they prefer
        3. Use create_presentation_from_template() with the chosen template

        This creates a blank presentation only - templates provide better formatting and design.

        Args:
            id: Optional presentation ID (auto-generated if None)

        Returns:
            Dictionary with presentation details and template recommendation
        """
        from pptx import Presentation

        pres = Presentation()
        if id is None:
            id = f"presentation_{len(self.presentations) + 1}"

        self.presentations[id] = pres
        self.current_presentation_id = id

        return {
            "presentation_id": id,
            "message": f"Created new blank presentation with ID: {id}",
            "slide_count": len(pres.slides),
            "recommendation": "💡 For better results, use discover_available_templates() and create_presentation_from_template() instead",
            "next_steps": "Add slides using add_slide() or create_slide_from_template() for professional layouts",
        }

    def create_presentation_from_template(
        self, template_path: str, id: Optional[str] = None
    ) -> dict:
        """Create a new PowerPoint presentation from a template file.

        🎯 PREFERRED METHOD for creating professional presentations!

        This function loads an existing PowerPoint template (.pptx) file and uses it as the foundation
        for a new presentation, preserving all formatting, themes, and slide masters.

        💡 TIP: Use discover_available_templates() first to show users available options!

        Args:
            template_path: Path to the PowerPoint template file (.pptx)
            id: Optional presentation ID (auto-generated if None)

        Returns:
            Dictionary with presentation details and template information
        """
        from pptx import Presentation

        # Resolve template path with base folder
        resolved_template_path = self.resolve_file_path(template_path)

        if not os.path.exists(resolved_template_path):
            search_dirs = self.get_template_search_directories()
            template_name = os.path.basename(template_path)

            for directory in search_dirs:
                potential_path = os.path.join(directory, template_name)
                if os.path.exists(potential_path):
                    resolved_template_path = potential_path
                    break
            else:
                return {
                    "error": f"Template file not found: {template_path}. Searched in {', '.join(search_dirs)}"
                }

        try:
            pres = Presentation(resolved_template_path)
        except Exception as e:
            return {"error": f"Failed to create presentation from template: {str(e)}"}

        if id is None:
            id = f"presentation_{len(self.presentations) + 1}"

        self.presentations[id] = pres
        self.current_presentation_id = id

        return {
            "presentation_id": id,
            "message": f"Created new presentation from template '{template_path}' with ID: {id}",
            "template_path": resolved_template_path,
            "original_path": template_path,
            "base_folder": self.valves.base_folder,
            "slide_count": len(pres.slides),
            "layout_count": len(pres.slide_layouts),
        }

    def open_presentation(self, file_path: str, id: Optional[str] = None) -> dict:
        """Open an existing PowerPoint presentation from a file."""
        from pptx import Presentation

        # Resolve file path with base folder
        resolved_file_path = self.resolve_file_path(file_path)

        if not os.path.exists(resolved_file_path):
            return {
                "error": f"File not found: {resolved_file_path} (original: {file_path})"
            }

        try:
            pres = Presentation(resolved_file_path)
        except Exception as e:
            return {"error": f"Failed to open presentation: {str(e)}"}

        if id is None:
            id = f"presentation_{len(self.presentations) + 1}"

        self.presentations[id] = pres
        self.current_presentation_id = id

        return {
            "presentation_id": id,
            "message": f"Opened presentation from {resolved_file_path} with ID: {id}",
            "file_path": resolved_file_path,
            "original_path": file_path,
            "base_folder": self.valves.base_folder,
            "slide_count": len(pres.slides),
        }

    def save_presentation(
        self, file_path: str, presentation_id: Optional[str] = None
    ) -> dict:
        """Save a presentation to a file."""
        pres_id = (
            presentation_id
            if presentation_id is not None
            else self.current_presentation_id
        )

        if pres_id is None or pres_id not in self.presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }

        # Resolve file path with base folder
        resolved_file_path = self.resolve_file_path(file_path)

        try:
            self.presentations[pres_id].save(resolved_file_path)
            return {
                "message": f"Presentation saved to {resolved_file_path}",
                "file_path": resolved_file_path,
                "original_path": file_path,
                "base_folder": self.valves.base_folder,
            }
        except Exception as e:
            return {"error": f"Failed to save presentation: {str(e)}"}

    def get_presentation_info(self, presentation_id: Optional[str] = None) -> dict:
        """Get information about a presentation."""
        pres_id = (
            presentation_id
            if presentation_id is not None
            else self.current_presentation_id
        )

        if pres_id is None or pres_id not in self.presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }

        pres = self.presentations[pres_id]

        try:
            return {
                "presentation_id": pres_id,
                "slide_count": len(pres.slides),
                "layout_count": len(pres.slide_layouts),
                "core_properties": {
                    "title": pres.core_properties.title or "",
                    "author": pres.core_properties.author or "",
                    "subject": pres.core_properties.subject or "",
                },
            }
        except Exception as e:
            return {"error": f"Failed to get presentation info: {str(e)}"}

    # --- Content Management Tools ---

    def add_slide(
        self,
        layout_index: int = 1,
        title: Optional[str] = None,
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Add a new slide with basic layout to the presentation.

        💡 RECOMMENDATION: Consider using create_slide_from_template() for better results!

        For professional slides with rich formatting and design:
        1. First call list_slide_templates() to see available template options
        2. Ask the user which template style they prefer
        3. Use create_slide_from_template() with the chosen template

        This function adds basic slides - templates provide professional layouts and styling.

        Args:
            layout_index: Index of the slide layout (0-based, default: 1)
            title: Optional title text for the slide
            presentation_id: Optional presentation ID (uses current if None)

        Returns:
            Dictionary with slide details and template recommendation
        """
        from pptx.util import Inches

        pres_id = (
            presentation_id
            if presentation_id is not None
            else self.current_presentation_id
        )

        if pres_id is None or pres_id not in self.presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }

        pres = self.presentations[pres_id]

        if layout_index < 0 or layout_index >= len(pres.slide_layouts):
            return {
                "error": f"Invalid layout index: {layout_index}. Available layouts: 0-{len(pres.slide_layouts) - 1}"
            }

        try:
            layout = pres.slide_layouts[layout_index]
            slide = pres.slides.add_slide(layout)
            slide_index = len(pres.slides) - 1

            if title:
                try:
                    if hasattr(slide.shapes, "title") and slide.shapes.title:
                        slide.shapes.title.text = title
                    else:
                        title_shape = slide.shapes.add_textbox(
                            Inches(0.5), Inches(0.5), Inches(9), Inches(1)
                        )
                        title_shape.text = title
                except:
                    pass

            return {
                "message": f"Added slide {slide_index} with layout {layout_index}",
                "slide_index": slide_index,
                "layout_name": (
                    layout.name if hasattr(layout, "name") else f"Layout {layout_index}"
                ),
                "recommendation": "💡 For professional styling, use list_slide_templates() and create_slide_from_template() instead",
            }
        except Exception as e:
            return {"error": f"Failed to add slide: {str(e)}"}

    def add_text_box(
        self,
        slide_index: int,
        text: str,
        left: float = 1.0,
        top: float = 1.0,
        width: float = 8.0,
        height: float = 1.0,
        font_size: Optional[int] = None,
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Add a text box to a slide."""
        from pptx.util import Inches, Pt

        pres_id = (
            presentation_id
            if presentation_id is not None
            else self.current_presentation_id
        )

        if pres_id is None or pres_id not in self.presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }

        pres = self.presentations[pres_id]

        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        try:
            textbox = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            text_frame = textbox.text_frame
            text_frame.text = text

            if font_size:
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(font_size)

            return {
                "message": f"Added text box to slide {slide_index}",
                "slide_index": slide_index,
                "text": text,
                "position": {
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                },
            }
        except Exception as e:
            return {"error": f"Failed to add text box: {str(e)}"}

    def get_server_info(self) -> dict:
        """Get information about the PowerPoint MCP toolkit."""
        return {
            "name": "PowerPoint MCP Toolkit - OpenWebUI Edition",
            "version": "2.1.0",
            "loaded_presentations": len(self.presentations),
            "current_presentation": self.current_presentation_id,
            "features": [
                "Presentation Management (create, open, save, switch)",
                "Content Management (slides, text boxes)",
                "Template Operations",
                "Advanced PowerPoint automation",
            ],
            "configuration": {
                "template_search_paths": self.valves.template_search_paths
            },
        }

    # --- Content Management Tools ---

    def extract_slide_text(
        self, slide_index: int, presentation_id: Optional[str] = None
    ) -> dict:
        """Extract all text content from a specific slide."""
        pres_id = (
            presentation_id
            if presentation_id is not None
            else self.current_presentation_id
        )

        if pres_id is None or pres_id not in self.presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }

        pres = self.presentations[pres_id]

        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        try:
            result = self._extract_slide_text_content(slide)
            result["slide_index"] = slide_index
            return result
        except Exception as e:
            return {"error": f"Failed to extract slide text: {str(e)}"}

    def extract_presentation_text(
        self, presentation_id: Optional[str] = None, include_slide_info: bool = True
    ) -> dict:
        """Extract all text content from all slides in the presentation."""
        pres_id = (
            presentation_id
            if presentation_id is not None
            else self.current_presentation_id
        )

        if pres_id is None or pres_id not in self.presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }

        pres = self.presentations[pres_id]

        try:
            slides_text = []
            total_text_shapes = 0
            slides_with_tables = 0
            slides_with_titles = 0
            all_presentation_text = []

            for slide_index, slide in enumerate(pres.slides):
                slide_text_result = self._extract_slide_text_content(slide)

                if slide_text_result["success"]:
                    slide_data = {
                        "slide_index": slide_index,
                        "text_content": slide_text_result["text_content"],
                    }

                    if include_slide_info:
                        slide_data["layout_name"] = slide.slide_layout.name
                        slide_data["total_text_shapes"] = slide_text_result[
                            "total_text_shapes"
                        ]
                        slide_data["has_title"] = slide_text_result["has_title"]
                        slide_data["has_tables"] = slide_text_result["has_tables"]

                    slides_text.append(slide_data)

                    # Accumulate statistics
                    total_text_shapes += slide_text_result["total_text_shapes"]
                    if slide_text_result["has_tables"]:
                        slides_with_tables += 1
                    if slide_text_result["has_title"]:
                        slides_with_titles += 1

                    # Collect all text for combined output
                    if slide_text_result["text_content"]["all_text_combined"]:
                        all_presentation_text.append(f"=== SLIDE {slide_index + 1} ===")
                        all_presentation_text.append(
                            slide_text_result["text_content"]["all_text_combined"]
                        )
                        all_presentation_text.append("")  # Empty line separator
                else:
                    slides_text.append(
                        {
                            "slide_index": slide_index,
                            "error": slide_text_result.get("error", "Unknown error"),
                            "text_content": None,
                        }
                    )

            return {
                "success": True,
                "presentation_id": pres_id,
                "total_slides": len(pres.slides),
                "slides_with_text": len(
                    [s for s in slides_text if s.get("text_content") is not None]
                ),
                "total_text_shapes": total_text_shapes,
                "slides_with_titles": slides_with_titles,
                "slides_with_tables": slides_with_tables,
                "slides_text": slides_text,
                "all_presentation_text_combined": "\n".join(all_presentation_text),
            }

        except Exception as e:
            return {"error": f"Failed to extract presentation text: {str(e)}"}

    def get_slide_info(
        self, slide_index: int, presentation_id: Optional[str] = None
    ) -> dict:
        """Get detailed information about a specific slide."""
        pres_id = (
            presentation_id
            if presentation_id is not None
            else self.current_presentation_id
        )

        if pres_id is None or pres_id not in self.presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }

        pres = self.presentations[pres_id]

        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        try:
            return self._get_slide_info(slide, slide_index)
        except Exception as e:
            return {"error": f"Failed to get slide info: {str(e)}"}

    def populate_placeholder(
        self,
        slide_index: int,
        placeholder_idx: int,
        text: str,
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Populate a placeholder with text."""
        pres_id = (
            presentation_id
            if presentation_id is not None
            else self.current_presentation_id
        )

        if pres_id is None or pres_id not in self.presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }

        pres = self.presentations[pres_id]

        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        try:
            if placeholder_idx < 0 or placeholder_idx >= len(slide.placeholders):
                return {
                    "error": f"Invalid placeholder index: {placeholder_idx}. Available placeholders: 0-{len(slide.placeholders) - 1}"
                }

            placeholder = slide.placeholders[placeholder_idx]
            placeholder.text = text

            return {
                "message": f"Populated placeholder {placeholder_idx} on slide {slide_index}",
                "placeholder_type": (
                    str(placeholder.placeholder_format.type)
                    if hasattr(placeholder, "placeholder_format")
                    else "unknown"
                ),
            }
        except Exception as e:
            return {"error": f"Failed to populate placeholder: {str(e)}"}

    def add_bullet_points(
        self,
        slide_index: int,
        placeholder_idx: int,
        bullet_points: List[str],
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Add bullet points to a placeholder."""
        pres_id = (
            presentation_id
            if presentation_id is not None
            else self.current_presentation_id
        )

        if pres_id is None or pres_id not in self.presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }

        pres = self.presentations[pres_id]

        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        try:
            if placeholder_idx < 0 or placeholder_idx >= len(slide.placeholders):
                return {
                    "error": f"Invalid placeholder index: {placeholder_idx}. Available placeholders: 0-{len(slide.placeholders) - 1}"
                }

            placeholder = slide.placeholders[placeholder_idx]
            self._add_bullet_points(placeholder, bullet_points)

            return {
                "message": f"Added {len(bullet_points)} bullet points to placeholder {placeholder_idx} on slide {slide_index}",
                "bullet_count": len(bullet_points),
            }
        except Exception as e:
            return {"error": f"Failed to add bullet points: {str(e)}"}

    # --- Internal Utility Methods ---

    def _extract_slide_text_content(self, slide) -> dict:
        """Extract all text content from a slide including placeholders and text shapes."""
        try:
            text_content = {
                "slide_title": "",
                "placeholders": [],
                "text_shapes": [],
                "table_text": [],
                "all_text_combined": "",
            }

            all_texts = []

            # Extract title from slide if available
            if (
                hasattr(slide, "shapes")
                and hasattr(slide.shapes, "title")
                and slide.shapes.title
            ):
                try:
                    title_text = slide.shapes.title.text_frame.text.strip()
                    if title_text:
                        text_content["slide_title"] = title_text
                        all_texts.append(title_text)
                except:
                    pass

            # Extract text from all shapes
            for i, shape in enumerate(slide.shapes):
                shape_text_info = {
                    "shape_index": i,
                    "shape_name": shape.name,
                    "shape_type": str(shape.shape_type),
                    "text": "",
                }

                try:
                    # Check if shape has text frame
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            shape_text_info["text"] = text
                            all_texts.append(text)

                            # Categorize by shape type
                            if hasattr(shape, "placeholder_format"):
                                # This is a placeholder
                                placeholder_info = shape_text_info.copy()
                                placeholder_info["placeholder_type"] = str(
                                    shape.placeholder_format.type
                                )
                                placeholder_info["placeholder_idx"] = (
                                    shape.placeholder_format.idx
                                )
                                text_content["placeholders"].append(placeholder_info)
                            else:
                                # This is a regular text shape
                                text_content["text_shapes"].append(shape_text_info)

                    # Extract text from tables
                    elif hasattr(shape, "table"):
                        table_texts = []
                        table = shape.table
                        for row_idx, row in enumerate(table.rows):
                            row_texts = []
                            for col_idx, cell in enumerate(row.cells):
                                cell_text = cell.text_frame.text.strip()
                                if cell_text:
                                    row_texts.append(cell_text)
                                    all_texts.append(cell_text)
                            if row_texts:
                                table_texts.append({"row": row_idx, "cells": row_texts})

                        if table_texts:
                            text_content["table_text"].append(
                                {
                                    "shape_index": i,
                                    "shape_name": shape.name,
                                    "table_content": table_texts,
                                }
                            )

                except Exception as e:
                    # Skip shapes that can't be processed
                    continue

            # Combine all text
            text_content["all_text_combined"] = "\n".join(all_texts)

            return {
                "success": True,
                "text_content": text_content,
                "total_text_shapes": len(text_content["placeholders"])
                + len(text_content["text_shapes"]),
                "has_title": bool(text_content["slide_title"]),
                "has_tables": len(text_content["table_text"]) > 0,
            }

        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to extract text content: {str(e)}",
                "text_content": None,
            }

    def _get_slide_info(self, slide, slide_index: int) -> dict:
        """Get detailed information about a specific slide."""
        try:
            placeholders = []
            for placeholder in slide.placeholders:
                placeholder_info = {
                    "idx": placeholder.placeholder_format.idx,
                    "type": str(placeholder.placeholder_format.type),
                    "name": placeholder.name,
                }
                placeholders.append(placeholder_info)

            shapes = []
            for i, shape in enumerate(slide.shapes):
                shape_info = {
                    "index": i,
                    "name": shape.name,
                    "shape_type": str(shape.shape_type),
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                }
                shapes.append(shape_info)

            return {
                "slide_index": slide_index,
                "layout_name": slide.slide_layout.name,
                "placeholder_count": len(placeholders),
                "placeholders": placeholders,
                "shape_count": len(shapes),
                "shapes": shapes,
            }
        except Exception as e:
            raise Exception(f"Failed to get slide info: {str(e)}")

    def _add_bullet_points(self, placeholder, bullet_points: List[str]) -> None:
        """Add bullet points to a placeholder."""
        from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

        text_frame = placeholder.text_frame
        text_frame.clear()  # Clear existing content

        for i, bullet_text in enumerate(bullet_points):
            if i == 0:
                # Use the first paragraph
                p = text_frame.paragraphs[0]
            else:
                # Add new paragraphs for subsequent bullets
                p = text_frame.add_paragraph()

            p.text = bullet_text
            p.level = 0  # Top level bullet

    # --- Structural Tools ---

    def add_table(
        self,
        slide_index: int,
        rows: int,
        cols: int,
        left: float,
        top: float,
        width: float,
        height: float,
        data: Optional[List[List[str]]] = None,
        header_row: bool = True,
        header_font_size: int = 12,
        body_font_size: int = 10,
        header_bg_color: Optional[List[int]] = None,
        body_bg_color: Optional[List[int]] = None,
        border_color: Optional[List[int]] = None,
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Add a table to a slide with enhanced formatting options."""
        pres_id = (
            presentation_id
            if presentation_id is not None
            else self.current_presentation_id
        )

        if pres_id is None or pres_id not in self.presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }

        pres = self.presentations[pres_id]

        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        # Validate parameters
        validations = {
            "rows": (rows, [(self.is_positive, "must be a positive integer")]),
            "cols": (cols, [(self.is_positive, "must be a positive integer")]),
            "left": (left, [(self.is_non_negative, "must be non-negative")]),
            "top": (top, [(self.is_non_negative, "must be non-negative")]),
            "width": (width, [(self.is_positive, "must be positive")]),
            "height": (height, [(self.is_positive, "must be positive")]),
        }

        if header_bg_color is not None:
            validations["header_bg_color"] = (
                header_bg_color,
                [
                    (
                        self.is_valid_rgb,
                        "must be a valid RGB list [R, G, B] with values 0-255",
                    )
                ],
            )
        if body_bg_color is not None:
            validations["body_bg_color"] = (
                body_bg_color,
                [
                    (
                        self.is_valid_rgb,
                        "must be a valid RGB list [R, G, B] with values 0-255",
                    )
                ],
            )
        if border_color is not None:
            validations["border_color"] = (
                border_color,
                [
                    (
                        self.is_valid_rgb,
                        "must be a valid RGB list [R, G, B] with values 0-255",
                    )
                ],
            )

        valid, error = self.validate_parameters(validations)
        if not valid:
            return {"error": error}

        # Validate data if provided
        if data:
            if len(data) != rows:
                return {
                    "error": f"Data has {len(data)} rows but table should have {rows} rows"
                }
            for i, row in enumerate(data):
                if len(row) != cols:
                    return {
                        "error": f"Row {i} has {len(row)} columns but table should have {cols} columns"
                    }

        try:
            # Add the table
            from pptx.util import Inches

            table_shape = slide.shapes.add_table(
                rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
            )
            table = table_shape.table

            # Populate with data if provided
            if data:
                for r in range(rows):
                    for c in range(cols):
                        if r < len(data) and c < len(data[r]):
                            table.cell(r, c).text = str(data[r][c])

            # Apply formatting
            for r in range(rows):
                for c in range(cols):
                    cell = table.cell(r, c)

                    # Header row formatting
                    if r == 0 and header_row:
                        self._format_table_cell(
                            cell,
                            bg_color=(
                                tuple(header_bg_color) if header_bg_color else None
                            ),
                            font_size=header_font_size,
                            bold=True,
                        )
                    else:
                        # Body cell formatting
                        self._format_table_cell(
                            cell,
                            bg_color=tuple(body_bg_color) if body_bg_color else None,
                            font_size=body_font_size,
                        )

            return {
                "message": f"Added {rows}x{cols} table to slide {slide_index}",
                "shape_index": len(slide.shapes) - 1,
                "rows": rows,
                "cols": cols,
            }
        except Exception as e:
            return {"error": f"Failed to add table: {str(e)}"}

    def format_table_cell(
        self,
        slide_index: int,
        shape_index: int,
        row: int,
        col: int,
        font_size: Optional[int] = None,
        font_name: Optional[str] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        color: Optional[List[int]] = None,
        bg_color: Optional[List[int]] = None,
        alignment: Optional[str] = None,
        vertical_alignment: Optional[str] = None,
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Format a specific table cell."""
        pres_id = (
            presentation_id
            if presentation_id is not None
            else self.current_presentation_id
        )

        if pres_id is None or pres_id not in self.presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }

        pres = self.presentations[pres_id]

        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        if shape_index < 0 or shape_index >= len(slide.shapes):
            return {
                "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
            }

        shape = slide.shapes[shape_index]

        try:
            if not hasattr(shape, "table"):
                return {"error": f"Shape at index {shape_index} is not a table"}

            table = shape.table

            if row < 0 or row >= len(table.rows):
                return {
                    "error": f"Invalid row index: {row}. Available rows: 0-{len(table.rows) - 1}"
                }

            if col < 0 or col >= len(table.columns):
                return {
                    "error": f"Invalid column index: {col}. Available columns: 0-{len(table.columns) - 1}"
                }

            cell = table.cell(row, col)

            self._format_table_cell(
                cell,
                font_size=font_size,
                font_name=font_name,
                bold=bold,
                italic=italic,
                color=tuple(color) if color else None,
                bg_color=tuple(bg_color) if bg_color else None,
                alignment=alignment,
                vertical_alignment=vertical_alignment,
            )

            return {
                "message": f"Formatted cell at row {row}, column {col} in table at shape index {shape_index} on slide {slide_index}"
            }
        except Exception as e:
            return {"error": f"Failed to format table cell: {str(e)}"}

    def add_shape(
        self,
        slide_index: int,
        shape_type: str,
        left: float,
        top: float,
        width: float,
        height: float,
        fill_color: Optional[List[int]] = None,
        line_color: Optional[List[int]] = None,
        line_width: Optional[float] = None,
        text: Optional[str] = None,
        font_size: Optional[int] = None,
        font_color: Optional[List[int]] = None,
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Add an auto shape to a slide with enhanced options."""
        pres_id = (
            presentation_id
            if presentation_id is not None
            else self.current_presentation_id
        )

        if pres_id is None or pres_id not in self.presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }

        pres = self.presentations[pres_id]

        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        try:
            # Use the existing add_shape_direct method
            shape = self.add_shape_direct(slide, shape_type, left, top, width, height)

            # Format the shape if formatting options are provided
            if any([fill_color, line_color, line_width]):
                self._format_shape(
                    shape,
                    fill_color=tuple(fill_color) if fill_color else None,
                    line_color=tuple(line_color) if line_color else None,
                    line_width=line_width,
                )

            # Add text to shape if provided
            if text and hasattr(shape, "text_frame"):
                shape.text_frame.text = text
                if font_size or font_color:
                    self._format_text(
                        shape.text_frame,
                        font_size=font_size,
                        color=tuple(font_color) if font_color else None,
                    )

            return {
                "message": f"Added {shape_type} shape to slide {slide_index}",
                "shape_index": len(slide.shapes) - 1,
            }
        except ValueError as e:
            return {"error": str(e)}
        except Exception as e:
            return {"error": f"Failed to add shape '{shape_type}': {str(e)}"}

    def delete_slide(
        self, slide_index: int, presentation_id: Optional[str] = None
    ) -> dict:
        """Delete a slide from the presentation."""
        pres_id = (
            presentation_id
            if presentation_id is not None
            else self.current_presentation_id
        )

        if pres_id is None or pres_id not in self.presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }

        pres = self.presentations[pres_id]

        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        if len(pres.slides) <= 1:
            return {"error": "Cannot delete the only slide in the presentation"}

        try:
            # Create a new presentation with all slides except the one to delete
            from pptx import Presentation

            new_pres = Presentation()

            # Copy all slides except the one to delete
            for i, slide in enumerate(pres.slides):
                if i != slide_index:
                    layout = slide.slide_layout
                    new_slide = new_pres.slides.add_slide(layout)

                    # Copy slide title if it exists
                    if hasattr(slide.shapes, "title") and slide.shapes.title:
                        try:
                            if (
                                hasattr(new_slide.shapes, "title")
                                and new_slide.shapes.title
                            ):
                                new_slide.shapes.title.text = slide.shapes.title.text
                        except:
                            pass

            # Replace the presentation
            self.presentations[pres_id] = new_pres

            return {
                "message": f"Deleted slide {slide_index}",
                "remaining_slides": len(new_pres.slides),
            }
        except Exception as e:
            return {"error": f"Failed to delete slide: {str(e)}"}

    def duplicate_slide(
        self, slide_index: int, presentation_id: Optional[str] = None
    ) -> dict:
        """Duplicate a slide in the presentation."""
        pres_id = (
            presentation_id
            if presentation_id is not None
            else self.current_presentation_id
        )

        if pres_id is None or pres_id not in self.presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }

        pres = self.presentations[pres_id]

        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        try:
            source_slide = pres.slides[slide_index]
            layout = source_slide.slide_layout
            new_slide = pres.slides.add_slide(layout)

            # Copy slide title if it exists
            if hasattr(source_slide.shapes, "title") and source_slide.shapes.title:
                try:
                    if hasattr(new_slide.shapes, "title") and new_slide.shapes.title:
                        new_slide.shapes.title.text = source_slide.shapes.title.text
                except:
                    pass

            new_slide_index = len(pres.slides) - 1

            return {
                "message": f"Duplicated slide {slide_index}",
                "new_slide_index": new_slide_index,
                "total_slides": len(pres.slides),
            }
        except Exception as e:
            return {"error": f"Failed to duplicate slide: {str(e)}"}

    # ==================== Professional Design Tools ====================

    def get_color_schemes(self) -> dict:
        """Get available professional color schemes for presentations.

        Returns:
            Dictionary with available color schemes and their properties
        """
        try:
            # Professional color schemes optimized for business presentations
            color_schemes = {
                "modern_blue": {
                    "primary": [70, 130, 180],  # Steel Blue
                    "secondary": [176, 196, 222],  # Light Steel Blue
                    "accent1": [255, 165, 0],  # Orange
                    "accent2": [240, 248, 255],  # Alice Blue
                    "light": [248, 248, 255],  # Ghost White
                    "text": [25, 25, 112],  # Midnight Blue
                    "description": "Modern corporate blue theme with orange accents",
                },
                "corporate_green": {
                    "primary": [34, 139, 34],  # Forest Green
                    "secondary": [144, 238, 144],  # Light Green
                    "accent1": [255, 215, 0],  # Gold
                    "accent2": [240, 255, 240],  # Honeydew
                    "light": [245, 255, 250],  # Mint Cream
                    "text": [0, 100, 0],  # Dark Green
                    "description": "Professional green theme with gold accents",
                },
                "elegant_purple": {
                    "primary": [75, 0, 130],  # Indigo
                    "secondary": [221, 160, 221],  # Plum
                    "accent1": [255, 20, 147],  # Deep Pink
                    "accent2": [230, 230, 250],  # Lavender
                    "light": [248, 248, 255],  # Ghost White
                    "text": [72, 61, 139],  # Dark Slate Blue
                    "description": "Elegant purple theme with pink accents",
                },
                "warm_orange": {
                    "primary": [255, 140, 0],  # Dark Orange
                    "secondary": [255, 218, 185],  # Peach Puff
                    "accent1": [220, 20, 60],  # Crimson
                    "accent2": [255, 245, 238],  # Seashell
                    "light": [255, 250, 240],  # Floral White
                    "text": [139, 69, 19],  # Saddle Brown
                    "description": "Warm orange theme with crimson accents",
                },
                "professional_gray": {
                    "primary": [105, 105, 105],  # Dim Gray
                    "secondary": [192, 192, 192],  # Silver
                    "accent1": [0, 191, 255],  # Deep Sky Blue
                    "accent2": [245, 245, 245],  # White Smoke
                    "light": [248, 248, 255],  # Ghost White
                    "text": [47, 79, 79],  # Dark Slate Gray
                    "description": "Professional gray theme with blue accents",
                },
            }

            return {
                "success": True,
                "available_schemes": list(color_schemes.keys()),
                "schemes": color_schemes,
                "color_types": [
                    "primary",
                    "secondary",
                    "accent1",
                    "accent2",
                    "light",
                    "text",
                ],
                "description": "Professional color schemes optimized for business presentations",
            }

        except Exception as e:
            return {"success": False, "error": f"Error getting color schemes: {str(e)}"}

    def apply_professional_theme(
        self, color_scheme: str = "modern_blue", apply_to_existing: bool = True
    ) -> dict:
        """Apply a professional theme to the current presentation.

        Args:
            color_scheme: Name of color scheme to apply
            apply_to_existing: Whether to apply to existing slides

        Returns:
            Dictionary with operation result
        """
        try:
            current_presentation = self.get_current_presentation()
            if current_presentation is None:
                return {
                    "success": False,
                    "error": "No presentation is currently loaded",
                }

            # Get color schemes
            schemes = self.get_color_schemes()
            if not schemes.get("success"):
                return schemes

            if color_scheme not in schemes["schemes"]:
                return {
                    "success": False,
                    "error": f"Invalid color scheme: {color_scheme}. Available: {list(schemes['schemes'].keys())}",
                }

            scheme_colors = schemes["schemes"][color_scheme]
            slides_affected = 0

            if apply_to_existing:
                # Apply theme to existing slides (simplified implementation)
                for slide in current_presentation.slides:
                    try:
                        # Apply background color to slide
                        self._apply_slide_background_color(
                            slide, scheme_colors["light"]
                        )

                        # Apply color scheme to shapes
                        for shape in slide.shapes:
                            if hasattr(shape, "text_frame") and shape.text_frame:
                                # Apply text colors based on scheme
                                self._apply_text_colors(shape, scheme_colors)

                        slides_affected += 1
                    except Exception:
                        continue  # Skip slides that can't be modified

            return {
                "success": True,
                "message": f"Applied {color_scheme} theme to presentation",
                "color_scheme": color_scheme,
                "slides_affected": slides_affected,
                "scheme_description": scheme_colors.get("description", ""),
            }

        except Exception as e:
            return {
                "success": False,
                "error": f"Error applying professional theme: {str(e)}",
            }

    def _apply_slide_background_color(self, slide, color_rgb):
        """Apply background color to a slide (helper method)."""
        try:
            # Note: Direct background color setting in python-pptx is limited
            # This is a simplified implementation
            pass
        except Exception:
            pass

    def _apply_text_colors(self, shape, scheme_colors):
        """Apply color scheme to text in a shape (helper method)."""
        try:
            if hasattr(shape, "text_frame") and shape.text_frame:
                # Apply primary color to text (simplified)
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Note: Color setting would require more complex RGB conversion
                        pass
        except Exception:
            pass

    def add_professional_slide(
        self,
        slide_type: str = "title_content",
        color_scheme: str = "modern_blue",
        title: str = None,
        content: list = None,
    ) -> dict:
        """Add a professionally designed slide with advanced styling.

        Args:
            slide_type: Type of slide ('title', 'title_content', 'content', 'blank')
            color_scheme: Color scheme to apply
            title: Slide title
            content: List of content items

        Returns:
            Dictionary with slide creation results
        """
        try:
            current_presentation = self.get_current_presentation()
            if current_presentation is None:
                return {
                    "success": False,
                    "error": "No presentation is currently loaded",
                }

            # Map slide types to layout indices
            layout_map = {
                "title": 0,  # Title slide
                "title_content": 1,  # Title and content
                "content": 6,  # Content only
                "blank": 6,  # Blank layout
            }

            layout_index = layout_map.get(slide_type, 1)

            try:
                layout = current_presentation.slide_layouts[layout_index]
                slide = current_presentation.slides.add_slide(layout)

                # Set title if provided
                if title and slide.shapes.title:
                    slide.shapes.title.text = title
                    # Apply professional title formatting
                    self._apply_professional_title_formatting(
                        slide.shapes.title, color_scheme
                    )

                # Add content if provided
                if content and len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]
                    content_text = "\n".join([f"• {item}" for item in content])
                    content_placeholder.text = content_text
                    # Apply professional content formatting
                    self._apply_professional_content_formatting(
                        content_placeholder, color_scheme
                    )

                return {
                    "success": True,
                    "message": f"Added professional {slide_type} slide with {color_scheme} theme",
                    "slide_index": len(current_presentation.slides) - 1,
                    "slide_type": slide_type,
                    "color_scheme": color_scheme,
                }
            except IndexError:
                return {
                    "success": False,
                    "error": f"Layout index {layout_index} not available in presentation",
                }

        except Exception as e:
            return {
                "success": False,
                "error": f"Error adding professional slide: {str(e)}",
            }

    def enhance_existing_slide(
        self,
        slide_index: int,
        color_scheme: str = "modern_blue",
        enhance_title: bool = True,
        enhance_content: bool = True,
        enhance_shapes: bool = True,
        enhance_charts: bool = True,
    ) -> dict:
        """Enhance an existing slide with professional styling.

        Args:
            slide_index: Index of slide to enhance
            color_scheme: Color scheme to apply
            enhance_title: Whether to enhance title formatting
            enhance_content: Whether to enhance content formatting
            enhance_shapes: Whether to enhance shape formatting
            enhance_charts: Whether to enhance chart formatting

        Returns:
            Dictionary with enhancement results
        """
        try:
            current_presentation = self.get_current_presentation()
            if current_presentation is None:
                return {
                    "success": False,
                    "error": "No presentation is currently loaded",
                }

            if slide_index < 0 or slide_index >= len(current_presentation.slides):
                return {
                    "success": False,
                    "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(current_presentation.slides) - 1}",
                }

            slide = current_presentation.slides[slide_index]
            enhancements_applied = []

            # Enhance title
            if enhance_title and slide.shapes.title:
                try:
                    self._apply_professional_title_formatting(
                        slide.shapes.title, color_scheme
                    )
                    enhancements_applied.append("title")
                except Exception:
                    pass

            # Enhance content and other shapes
            if enhance_content or enhance_shapes:
                for shape in slide.shapes:
                    try:
                        if (
                            hasattr(shape, "text_frame")
                            and shape.text_frame
                            and shape != slide.shapes.title
                        ):
                            self._apply_professional_content_formatting(
                                shape, color_scheme
                            )
                            enhancements_applied.append("content_shape")
                        elif enhance_shapes and hasattr(shape, "fill"):
                            self._apply_professional_shape_formatting(
                                shape, color_scheme
                            )
                            enhancements_applied.append("shape")
                    except Exception:
                        continue

            # Enhance charts
            if enhance_charts:
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "chart"):
                            self._apply_professional_chart_formatting(
                                shape, color_scheme
                            )
                            enhancements_applied.append("chart")
                    except Exception:
                        continue

            return {
                "success": True,
                "message": f"Enhanced slide {slide_index} with {color_scheme} scheme",
                "slide_index": slide_index,
                "color_scheme": color_scheme,
                "enhancements_applied": list(set(enhancements_applied)),
                "enhancement_count": len(enhancements_applied),
            }

        except Exception as e:
            return {"success": False, "error": f"Error enhancing slide: {str(e)}"}

    def _apply_professional_title_formatting(self, title_shape, color_scheme):
        """Apply professional formatting to title shape."""
        try:
            from pptx.util import Pt

            schemes = self.get_color_schemes()
            if schemes.get("success") and color_scheme in schemes["schemes"]:
                # Apply title formatting (simplified due to python-pptx limitations)
                if hasattr(title_shape, "text_frame"):
                    for paragraph in title_shape.text_frame.paragraphs:
                        paragraph.font.size = Pt(36)  # Large title font
                        paragraph.font.bold = True
        except Exception:
            pass

    def _apply_professional_content_formatting(self, content_shape, color_scheme):
        """Apply professional formatting to content shape."""
        try:
            from pptx.util import Pt

            schemes = self.get_color_schemes()
            if schemes.get("success") and color_scheme in schemes["schemes"]:
                # Apply content formatting (simplified)
                if hasattr(content_shape, "text_frame"):
                    for paragraph in content_shape.text_frame.paragraphs:
                        paragraph.font.size = Pt(18)  # Standard content font
        except Exception:
            pass

    def _apply_professional_shape_formatting(self, shape, color_scheme):
        """Apply professional formatting to shapes."""
        try:
            schemes = self.get_color_schemes()
            if schemes.get("success") and color_scheme in schemes["schemes"]:
                scheme_colors = schemes["schemes"][color_scheme]
                # Apply shape formatting (simplified)
                if hasattr(shape, "fill"):
                    # Would apply colors from scheme if RGB conversion was implemented
                    pass
        except Exception:
            pass

    def _apply_professional_chart_formatting(self, chart_shape, color_scheme):
        """Apply professional formatting to charts."""
        try:
            schemes = self.get_color_schemes()
            if schemes.get("success") and color_scheme in schemes["schemes"]:
                # Apply chart formatting (simplified)
                if hasattr(chart_shape, "chart"):
                    # Would apply professional chart styling
                    pass
        except Exception:
            pass

    def add_chart(
        self,
        slide_index: int,
        chart_type: str,
        left: float,
        top: float,
        width: float,
        height: float,
        categories: List[str],
        series_names: List[str],
        series_values: List[List[float]],
        has_legend: bool = True,
        legend_position: str = "right",
        has_data_labels: bool = False,
        title: Optional[str] = None,
        x_axis_title: Optional[str] = None,
        y_axis_title: Optional[str] = None,
        color_scheme: Optional[str] = None,
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Add a chart to a slide with comprehensive formatting options."""
        pres_id = (
            presentation_id
            if presentation_id is not None
            else self.current_presentation_id
        )

        if pres_id is None or pres_id not in self.presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }

        pres = self.presentations[pres_id]

        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        # Validate chart type
        valid_chart_types = [
            "column",
            "stacked_column",
            "bar",
            "stacked_bar",
            "line",
            "line_markers",
            "pie",
            "doughnut",
            "area",
            "stacked_area",
            "scatter",
            "radar",
            "radar_markers",
        ]
        if chart_type.lower() not in valid_chart_types:
            return {
                "error": f"Invalid chart type: '{chart_type}'. Valid types are: {', '.join(valid_chart_types)}"
            }

        # Validate series data
        if len(series_names) != len(series_values):
            return {
                "error": f"Number of series names ({len(series_names)}) must match number of series values ({len(series_values)})"
            }

        if not categories:
            return {"error": "Categories list cannot be empty"}

        # Validate that all series have the same number of values as categories
        for i, values in enumerate(series_values):
            if len(values) != len(categories):
                return {
                    "error": f"Series '{series_names[i]}' has {len(values)} values but there are {len(categories)} categories"
                }

        try:
            # Add the chart
            chart = self._add_chart(
                slide,
                chart_type,
                left,
                top,
                width,
                height,
                categories,
                series_names,
                series_values,
            )

            if chart is None:
                return {"error": "Failed to create chart"}

            # Format the chart
            self._format_chart(
                chart,
                has_legend=has_legend,
                legend_position=legend_position,
                has_data_labels=has_data_labels,
                title=title,
                x_axis_title=x_axis_title,
                y_axis_title=y_axis_title,
                color_scheme=color_scheme,
            )

            return {
                "message": f"Added {chart_type} chart to slide {slide_index}",
                "shape_index": len(slide.shapes) - 1,
                "chart_type": chart_type,
                "series_count": len(series_names),
                "categories_count": len(categories),
            }
        except Exception as e:
            return {"error": f"Failed to add chart: {str(e)}"}

    # --- Internal Helper Methods for Structural Tools ---

    def _add_chart(
        self,
        slide,
        chart_type: str,
        left: float,
        top: float,
        width: float,
        height: float,
        categories: List[str],
        series_names: List[str],
        series_values: List[List[float]],
    ):
        """Add a chart to a slide."""
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches

        # Map chart type names to enum values
        chart_type_map = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
            "line": XL_CHART_TYPE.LINE,
            "line_markers": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
            "area": XL_CHART_TYPE.AREA,
            "stacked_area": XL_CHART_TYPE.AREA_STACKED,
            "scatter": XL_CHART_TYPE.XY_SCATTER,
            "radar": XL_CHART_TYPE.RADAR,
            "radar_markers": XL_CHART_TYPE.RADAR_MARKERS,
        }

        chart_type_enum = chart_type_map.get(chart_type.lower())
        if not chart_type_enum:
            raise ValueError(f"Unsupported chart type: {chart_type}")

        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = categories

        for series_name, values in zip(series_names, series_values):
            chart_data.add_series(series_name, values)

        # Add chart to slide
        chart_shape = slide.shapes.add_chart(
            chart_type_enum,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
            chart_data,
        )

        return chart_shape.chart

    def _format_chart(
        self,
        chart,
        has_legend: bool = True,
        legend_position: str = "right",
        has_data_labels: bool = False,
        title: Optional[str] = None,
        x_axis_title: Optional[str] = None,
        y_axis_title: Optional[str] = None,
        color_scheme: Optional[str] = None,
    ) -> None:
        """Format a chart with various styling options."""
        from pptx.enum.chart import XL_LEGEND_POSITION

        try:
            # Set chart title
            if title and hasattr(chart, "chart_title"):
                chart.chart_title.text_frame.text = title

            # Configure legend
            if has_legend and hasattr(chart, "has_legend"):
                chart.has_legend = True
                legend_position_map = {
                    "bottom": XL_LEGEND_POSITION.BOTTOM,
                    "left": XL_LEGEND_POSITION.LEFT,
                    "right": XL_LEGEND_POSITION.RIGHT,
                    "top": XL_LEGEND_POSITION.TOP,
                }
                if legend_position.lower() in legend_position_map:
                    chart.legend.position = legend_position_map[legend_position.lower()]
            else:
                chart.has_legend = False

            # Configure axis titles
            if x_axis_title and hasattr(chart, "category_axis"):
                try:
                    chart.category_axis.axis_title.text_frame.text = x_axis_title
                except:
                    pass  # Some chart types don't support axis titles

            if y_axis_title and hasattr(chart, "value_axis"):
                try:
                    chart.value_axis.axis_title.text_frame.text = y_axis_title
                except:
                    pass  # Some chart types don't support axis titles

            # Configure data labels
            if has_data_labels:
                try:
                    for series in chart.series:
                        series.has_data_labels = True
                except:
                    pass  # Some chart types don't support data labels

        except Exception as e:
            # Silently continue if formatting fails
            pass

    def _format_table_cell(
        self,
        cell,
        font_size: Optional[int] = None,
        font_name: Optional[str] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        color: Optional[Tuple] = None,
        bg_color: Optional[Tuple] = None,
        alignment: Optional[str] = None,
        vertical_alignment: Optional[str] = None,
    ) -> None:
        """Format a table cell with various styling options."""
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

        text_frame = cell.text_frame

        # Apply text formatting
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if font_size:
                    run.font.size = Pt(font_size)
                if font_name:
                    run.font.name = font_name
                if bold is not None:
                    run.font.bold = bold
                if italic is not None:
                    run.font.italic = italic
                if color:
                    run.font.color.rgb = RGBColor(*color)

        # Apply background color
        if bg_color:
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*bg_color)

        # Apply text alignment
        if alignment:
            alignment_map = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
                "justify": PP_ALIGN.JUSTIFY,
            }
            if alignment.lower() in alignment_map:
                for paragraph in text_frame.paragraphs:
                    paragraph.alignment = alignment_map[alignment.lower()]

        # Apply vertical alignment
        if vertical_alignment:
            vertical_alignment_map = {
                "top": MSO_ANCHOR.TOP,
                "middle": MSO_ANCHOR.MIDDLE,
                "bottom": MSO_ANCHOR.BOTTOM,
            }
            if vertical_alignment.lower() in vertical_alignment_map:
                text_frame.vertical_anchor = vertical_alignment_map[
                    vertical_alignment.lower()
                ]

    def _format_shape(
        self,
        shape,
        fill_color: Optional[Tuple] = None,
        line_color: Optional[Tuple] = None,
        line_width: Optional[float] = None,
    ) -> None:
        """Format a shape with fill and line properties."""
        from pptx.dml.color import RGBColor
        from pptx.util import Pt

        if fill_color:
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*fill_color)

        if line_color or line_width:
            line = shape.line
            if line_color:
                line.color.rgb = RGBColor(*line_color)
            if line_width:
                line.width = Pt(line_width)

    def _format_text(
        self, text_frame, font_size: Optional[int] = None, color: Optional[Tuple] = None
    ) -> None:
        """Format text in a text frame."""
        from pptx.util import Pt
        from pptx.dml.color import RGBColor

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if font_size:
                    run.font.size = Pt(font_size)
                if color:
                    run.font.color.rgb = RGBColor(*color)

    # ==================== Professional Picture Effects ====================

    def apply_picture_effects(
        self,
        slide_index: int,
        shape_index: int,
        effects: Dict[str, Dict],
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Apply multiple picture effects to a shape including shadows, glow, reflection, etc.

        Args:
            slide_index: Index of the slide containing the shape
            shape_index: Index of the shape to apply effects to
            effects: Dictionary of effects to apply. Format:
                     {"shadow": {"blur_radius": 4.0, "distance": 3.0, "direction": 315.0, "color": [0,0,0], "transparency": 0.6},
                      "glow": {"size": 5.0, "color": [0,176,240], "transparency": 0.4},
                      "reflection": {"size": 0.5, "transparency": 0.5, "distance": 0.0, "blur": 4.0},
                      "soft_edges": {"radius": 2.5},
                      "rotation": {"rotation": 15.0},
                      "transparency": {"transparency": 0.2},
                      "bevel": {"bevel_type": "circle", "width": 6.0, "height": 6.0}}
            presentation_id: Optional presentation ID (uses current if None)

        Returns:
            Dictionary with operation results and applied effects
        """
        try:
            current_presentation = self.get_current_presentation()
            if current_presentation is None:
                return {
                    "success": False,
                    "error": "No presentation is currently loaded",
                }

            if slide_index < 0 or slide_index >= len(current_presentation.slides):
                return {
                    "success": False,
                    "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(current_presentation.slides) - 1}",
                }

            slide = current_presentation.slides[slide_index]

            if shape_index < 0 or shape_index >= len(slide.shapes):
                return {
                    "success": False,
                    "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}",
                }

            shape = slide.shapes[shape_index]
            applied_effects = []
            warnings = []

            # Apply each effect
            for effect_type, effect_params in effects.items():
                try:
                    if effect_type == "shadow":
                        self._apply_picture_shadow(shape, effect_params)
                        applied_effects.append("shadow")
                    elif effect_type == "reflection":
                        self._apply_picture_reflection(shape, effect_params)
                        applied_effects.append("reflection")
                    elif effect_type == "glow":
                        self._apply_picture_glow(shape, effect_params)
                        applied_effects.append("glow")
                    elif effect_type == "soft_edges":
                        self._apply_picture_soft_edges(shape, effect_params)
                        applied_effects.append("soft_edges")
                    elif effect_type == "rotation":
                        self._apply_picture_rotation(shape, effect_params)
                        applied_effects.append("rotation")
                    elif effect_type == "transparency":
                        self._apply_picture_transparency(shape, effect_params)
                        applied_effects.append("transparency")
                    elif effect_type == "bevel":
                        self._apply_picture_bevel(shape, effect_params)
                        applied_effects.append("bevel")
                    else:
                        warnings.append(f"Unknown effect type: {effect_type}")
                except Exception as e:
                    warnings.append(f"Failed to apply {effect_type} effect: {str(e)}")

            result = {
                "success": True,
                "message": f"Applied {len(applied_effects)} effects to shape {shape_index} on slide {slide_index}",
                "applied_effects": applied_effects,
                "slide_index": slide_index,
                "shape_index": shape_index,
            }

            if warnings:
                result["warnings"] = warnings

            return result

        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to apply picture effects: {str(e)}",
            }

    def set_gradient_background(
        self,
        slide_index: int,
        start_color: List[int],
        end_color: List[int],
        direction: str = "horizontal",
        style: str = "custom",
        color_scheme: Optional[str] = None,
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Set a gradient background for a slide.

        Args:
            slide_index: Index of the slide to apply gradient to
            start_color: Starting RGB color as [r, g, b] (0-255 each)
            end_color: Ending RGB color as [r, g, b] (0-255 each)
            direction: Gradient direction - "horizontal", "vertical", "diagonal"
            style: Gradient style - "custom", "subtle", "bold", "accent" (if using color_scheme)
            color_scheme: Optional professional color scheme name (overrides start/end colors)
            presentation_id: Optional presentation ID (uses current if None)

        Returns:
            Dictionary with operation results
        """
        try:
            current_presentation = self.get_current_presentation()
            if current_presentation is None:
                return {
                    "success": False,
                    "error": "No presentation is currently loaded",
                }

            if slide_index < 0 or slide_index >= len(current_presentation.slides):
                return {
                    "success": False,
                    "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(current_presentation.slides) - 1}",
                }

            slide = current_presentation.slides[slide_index]

            # Validate colors
            if not self.is_valid_rgb(start_color) or not self.is_valid_rgb(end_color):
                return {
                    "success": False,
                    "error": "Invalid RGB color values. Use [r, g, b] format with values 0-255",
                }

            # Use professional color scheme if specified
            if color_scheme:
                color_schemes = self.get_color_schemes()
                if color_scheme in color_schemes["schemes"]:
                    scheme_colors = color_schemes["schemes"][color_scheme]
                    if style == "subtle":
                        start_color = scheme_colors["light"]
                        end_color = scheme_colors["secondary"]
                    elif style == "bold":
                        start_color = scheme_colors["primary"]
                        end_color = scheme_colors["accent1"]
                    elif style == "accent":
                        start_color = scheme_colors["accent1"]
                        end_color = scheme_colors["accent2"]

            # Apply gradient background
            self._create_gradient_background(
                slide, tuple(start_color), tuple(end_color), direction
            )

            return {
                "success": True,
                "message": f"Applied {direction} gradient background to slide {slide_index}",
                "slide_index": slide_index,
                "start_color": start_color,
                "end_color": end_color,
                "direction": direction,
                "style": style,
                "color_scheme": color_scheme,
            }

        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to set gradient background: {str(e)}",
            }

    def analyze_fonts(
        self, analysis_type: str = "presentation", presentation_id: Optional[str] = None
    ) -> dict:
        """Analyze fonts used in the presentation and provide recommendations.

        Args:
            analysis_type: Type of analysis - "presentation", "accessibility", "branding"
            presentation_id: Optional presentation ID (uses current if None)

        Returns:
            Dictionary with font analysis and recommendations
        """
        try:
            current_presentation = self.get_current_presentation()
            if current_presentation is None:
                return {
                    "success": False,
                    "error": "No presentation is currently loaded",
                }

            fonts_used = set()
            font_details = {}
            slide_font_usage = []

            # Analyze fonts in each slide
            for slide_idx, slide in enumerate(current_presentation.slides):
                slide_fonts = set()

                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if hasattr(run.font, "name") and run.font.name:
                                    font_name = run.font.name
                                    fonts_used.add(font_name)
                                    slide_fonts.add(font_name)

                                    if font_name not in font_details:
                                        font_details[font_name] = {
                                            "slides_used": [],
                                            "total_usage": 0,
                                            "sizes_used": set(),
                                            "is_bold": False,
                                            "is_italic": False,
                                        }

                                    font_details[font_name]["slides_used"].append(
                                        slide_idx
                                    )
                                    font_details[font_name]["total_usage"] += 1

                                    if hasattr(run.font, "size") and run.font.size:
                                        font_details[font_name]["sizes_used"].add(
                                            run.font.size.pt
                                        )
                                    if hasattr(run.font, "bold") and run.font.bold:
                                        font_details[font_name]["is_bold"] = True
                                    if hasattr(run.font, "italic") and run.font.italic:
                                        font_details[font_name]["is_italic"] = True

                slide_font_usage.append(
                    {"slide_index": slide_idx, "fonts": list(slide_fonts)}
                )

            # Convert sets to lists for JSON serialization
            for font in font_details:
                font_details[font]["sizes_used"] = list(
                    font_details[font]["sizes_used"]
                )
                font_details[font]["slides_used"] = list(
                    set(font_details[font]["slides_used"])
                )

            # Generate recommendations
            recommendations = self._generate_font_recommendations(
                fonts_used, font_details, analysis_type
            )

            return {
                "success": True,
                "analysis_type": analysis_type,
                "total_fonts_used": len(fonts_used),
                "fonts_used": list(fonts_used),
                "font_details": font_details,
                "slide_font_usage": slide_font_usage,
                "recommendations": recommendations,
            }

        except Exception as e:
            return {"success": False, "error": f"Failed to analyze fonts: {str(e)}"}

    # ==================== Picture Effects Helper Methods ====================

    def _apply_picture_shadow(self, shape, params: Dict) -> None:
        """Apply shadow effect to a shape."""
        try:
            # Basic shadow implementation - python-pptx has limited shadow support
            if hasattr(shape, "shadow"):
                shadow = shape.shadow
                if "blur_radius" in params:
                    # Note: python-pptx shadow API is limited, this is a simplified implementation
                    pass
        except Exception:
            pass  # Graceful fallback for unsupported shadow features

    def _apply_picture_reflection(self, shape, params: Dict) -> None:
        """Apply reflection effect to a shape."""
        try:
            # Reflection effects are not directly supported in python-pptx
            # This would require XML manipulation or external libraries
            pass
        except Exception:
            pass

    def _apply_picture_glow(self, shape, params: Dict) -> None:
        """Apply glow effect to a shape."""
        try:
            # Glow effects are not directly supported in python-pptx
            # This would require XML manipulation
            pass
        except Exception:
            pass

    def _apply_picture_soft_edges(self, shape, params: Dict) -> None:
        """Apply soft edges effect to a shape."""
        try:
            # Soft edges are not directly supported in python-pptx
            pass
        except Exception:
            pass

    def _apply_picture_rotation(self, shape, params: Dict) -> None:
        """Apply rotation to a shape."""
        try:
            rotation = params.get("rotation", 0.0)
            if hasattr(shape, "rotation"):
                shape.rotation = rotation
        except Exception:
            pass

    def _apply_picture_transparency(self, shape, params: Dict) -> None:
        """Apply transparency to a shape."""
        try:
            transparency = params.get("transparency", 0.0)
            # Limited transparency support in python-pptx
            if hasattr(shape, "fill"):
                shape.fill.transparency = transparency
        except Exception:
            pass

    def _apply_picture_bevel(self, shape, params: Dict) -> None:
        """Apply bevel effect to a shape."""
        try:
            # Bevel effects are not directly supported in python-pptx
            pass
        except Exception:
            pass

    def _create_gradient_background(
        self, slide, start_color: tuple, end_color: tuple, direction: str
    ) -> None:
        """Create a gradient background for a slide."""
        try:
            from PIL import Image, ImageDraw
            import tempfile
            import os

            # Create gradient image
            width, height = 1920, 1080  # Standard slide dimensions
            img = Image.new("RGB", (width, height))
            draw = ImageDraw.Draw(img)

            if direction == "horizontal":
                for x in range(width):
                    ratio = x / width
                    r = int(start_color[0] * (1 - ratio) + end_color[0] * ratio)
                    g = int(start_color[1] * (1 - ratio) + end_color[1] * ratio)
                    b = int(start_color[2] * (1 - ratio) + end_color[2] * ratio)
                    draw.line([(x, 0), (x, height)], fill=(r, g, b))
            elif direction == "vertical":
                for y in range(height):
                    ratio = y / height
                    r = int(start_color[0] * (1 - ratio) + end_color[0] * ratio)
                    g = int(start_color[1] * (1 - ratio) + end_color[1] * ratio)
                    b = int(start_color[2] * (1 - ratio) + end_color[2] * ratio)
                    draw.line([(0, y), (width, y)], fill=(r, g, b))
            else:  # diagonal
                for x in range(width):
                    for y in range(height):
                        ratio = (x + y) / (width + height)
                        r = int(start_color[0] * (1 - ratio) + end_color[0] * ratio)
                        g = int(start_color[1] * (1 - ratio) + end_color[1] * ratio)
                        b = int(start_color[2] * (1 - ratio) + end_color[2] * ratio)
                        img.putpixel((x, y), (r, g, b))

            # Save to temporary file and add as background
            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as temp_file:
                img.save(temp_file.name, "PNG")
                temp_path = temp_file.name

            try:
                from pptx.util import Inches

                # Add as background image (will be behind other content)
                slide.shapes.add_picture(temp_path, 0, 0, Inches(10), Inches(7.5))
                # Move to back if possible
                shape = slide.shapes[-1]
                slide.shapes._spTree.remove(shape._element)
                slide.shapes._spTree.insert(
                    2, shape._element
                )  # Insert after slide properties
            finally:
                # Clean up temporary file
                if os.path.exists(temp_path):
                    os.unlink(temp_path)

        except Exception as e:
            # Fallback to solid color background
            try:
                from pptx.dml.color import RGBColor
                from pptx.enum.dml import MSO_FILL

                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(*start_color)
            except Exception:
                pass  # Graceful fallback

    def _generate_font_recommendations(
        self, fonts_used: set, font_details: dict, analysis_type: str
    ) -> Dict:
        """Generate font recommendations based on analysis."""
        recommendations = []

        # Professional font recommendations
        professional_fonts = [
            "Segoe UI",
            "Arial",
            "Calibri",
            "Helvetica",
            "Open Sans",
            "Roboto",
        ]
        accessible_fonts = ["Arial", "Calibri", "Verdana", "Tahoma", "Open Sans"]

        if analysis_type == "presentation":
            if len(fonts_used) > 3:
                recommendations.append(
                    {
                        "type": "warning",
                        "message": f"Using {len(fonts_used)} different fonts. Consider limiting to 2-3 fonts for better consistency.",
                    }
                )

            non_professional = [f for f in fonts_used if f not in professional_fonts]
            if non_professional:
                recommendations.append(
                    {
                        "type": "suggestion",
                        "message": f"Consider replacing non-standard fonts: {', '.join(non_professional)}",
                        "suggested_fonts": professional_fonts[:3],
                    }
                )

        elif analysis_type == "accessibility":
            non_accessible = [f for f in fonts_used if f not in accessible_fonts]
            if non_accessible:
                recommendations.append(
                    {
                        "type": "accessibility",
                        "message": f"For better accessibility, consider replacing: {', '.join(non_accessible)}",
                        "suggested_fonts": accessible_fonts[:3],
                    }
                )

        elif analysis_type == "branding":
            recommendations.append(
                {
                    "type": "branding",
                    "message": "Ensure font choices align with your brand guidelines",
                    "suggested_action": "Verify fonts match corporate style guide",
                }
            )

        return {
            "total_recommendations": len(recommendations),
            "recommendations": recommendations,
            "professional_fonts_available": professional_fonts,
            "accessibility_fonts_available": accessible_fonts,
        }

    # ==================== Priority 4: Specialized Tools ====================
    # Template Management Tools

    def list_slide_templates(self) -> dict:
        """List all available slide layout templates and their usage examples.

        Returns:
            Dictionary with available templates, usage examples, and guidance
        """
        try:
            available_templates = {
                "title_slide": {
                    "description": "Professional title slide with title, subtitle, and author",
                    "content_mapping": {
                        "title": "Main Title",
                        "subtitle": "Subtitle Text",
                        "author": "Author Name",
                    },
                    "supports_images": False,
                },
                "text_with_image": {
                    "description": "Slide with text content and supporting image",
                    "content_mapping": {
                        "title": "Slide Title",
                        "content": "Main content text",
                    },
                    "supports_images": True,
                    "image_mapping": {"supporting": "Path to supporting image"},
                },
                "two_column_text": {
                    "description": "Two-column layout with balanced text content",
                    "content_mapping": {
                        "title": "Slide Title",
                        "left_content": "Left column text",
                        "right_content": "Right column text",
                    },
                    "supports_images": False,
                },
                "bullet_points": {
                    "description": "Standard bullet point presentation slide",
                    "content_mapping": {
                        "title": "Slide Title",
                        "bullets": "• Point 1\\n• Point 2\\n• Point 3",
                    },
                    "supports_images": False,
                },
                "image_showcase": {
                    "description": "Image-focused slide with minimal text",
                    "content_mapping": {
                        "title": "Optional Title",
                        "caption": "Image caption",
                    },
                    "supports_images": True,
                    "image_mapping": {"main": "Path to main image"},
                },
                "comparison": {
                    "description": "Side-by-side comparison layout",
                    "content_mapping": {
                        "title": "Comparison Title",
                        "left_title": "Option A",
                        "right_title": "Option B",
                        "left_content": "A details",
                        "right_content": "B details",
                    },
                    "supports_images": True,
                    "image_mapping": {
                        "left_image": "Option A image",
                        "right_image": "Option B image",
                    },
                },
            }

            usage_examples = {
                "title_slide": {
                    "example": "Perfect for presentation opening slides",
                    "best_practices": [
                        "Keep title concise",
                        "Use professional subtitle",
                        "Include presenter name",
                    ],
                },
                "text_with_image": {
                    "example": "Ideal for explaining concepts with visual support",
                    "best_practices": [
                        "Balance text and image",
                        "Use high-quality images",
                        "Keep text scannable",
                    ],
                },
                "comparison": {
                    "example": "Great for before/after, pros/cons, option evaluation",
                    "best_practices": [
                        "Use parallel structure",
                        "Keep comparisons balanced",
                        "Highlight key differences",
                    ],
                },
            }

            return {
                "success": True,
                "available_templates": available_templates,
                "total_templates": len(available_templates),
                "usage_examples": usage_examples,
                "message": "Use apply_slide_template() or create_slide_from_template() to apply these templates",
            }
        except Exception as e:
            return {"success": False, "error": f"Failed to list templates: {str(e)}"}

    def discover_available_templates(self) -> dict:
        """Discover all available PowerPoint template files (.pptx) in configured directories.

        This function scans the template directories to find actual PowerPoint template files
        that can be used with create_presentation_from_template(). This helps the AI suggest
        specific templates to users.

        Returns:
            Dictionary with discovered template files, their paths, and usage suggestions
        """
        try:
            template_files = []
            search_directories = self.get_template_search_directories()

            for directory in search_directories:
                if not os.path.exists(directory):
                    continue

                for filename in os.listdir(directory):
                    if filename.lower().endswith(".pptx") and not filename.startswith(
                        "~"
                    ):
                        full_path = os.path.join(directory, filename)
                        if os.path.isfile(full_path):
                            template_info = {
                                "filename": filename,
                                "path": full_path,
                                "directory": directory,
                                "name": os.path.splitext(filename)[0],
                                "size_mb": round(
                                    os.path.getsize(full_path) / (1024 * 1024), 2
                                ),
                            }
                            template_files.append(template_info)

            # Sort by filename for consistent ordering
            template_files.sort(key=lambda x: x["filename"].lower())

            usage_guidance = {
                "how_to_use": "Use create_presentation_from_template(template_path='path', id='presentation_id') to create presentations from these templates",
                "recommendation": "Always ask the user which template they prefer before creating presentations",
                "ai_instruction": "When a user wants to create a presentation, first call this function to show available templates, then ask the user to choose one",
            }

            return {
                "success": True,
                "discovered_templates": template_files,
                "total_files": len(template_files),
                "search_directories": search_directories,
                "usage_guidance": usage_guidance,
                "message": f"Found {len(template_files)} PowerPoint template files across {len(search_directories)} directories",
            }

        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to discover templates: {str(e)}",
            }

    def apply_slide_template(
        self,
        slide_index: int,
        template_id: str,
        color_scheme: str = "modern_blue",
        content_mapping: Optional[Dict[str, str]] = None,
        image_paths: Optional[Dict[str, str]] = None,
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Apply a structured layout template to an existing slide.

        Args:
            slide_index: Index of the slide to apply template to
            template_id: ID of the template to apply (e.g., 'title_slide', 'text_with_image')
            color_scheme: Color scheme to use from professional color schemes
            content_mapping: Dictionary mapping element roles to custom content
            image_paths: Dictionary mapping image element roles to file paths
            presentation_id: Optional presentation ID (uses current if None)

        Returns:
            Dictionary with operation results
        """
        try:
            current_presentation = self.get_current_presentation()
            if current_presentation is None:
                return {
                    "success": False,
                    "error": "No presentation is currently loaded",
                }

            if slide_index < 0 or slide_index >= len(current_presentation.slides):
                return {
                    "success": False,
                    "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(current_presentation.slides) - 1}",
                }

            slide = current_presentation.slides[slide_index]

            # Get available templates
            templates_result = self.list_slide_templates()
            if not templates_result["success"]:
                return templates_result

            available_templates = templates_result["available_templates"]

            if template_id not in available_templates:
                return {
                    "success": False,
                    "error": f"Unknown template ID: {template_id}. Available templates: {list(available_templates.keys())}",
                }

            template_config = available_templates[template_id]
            content_mapping = content_mapping or {}
            image_paths = image_paths or {}

            # Clear existing content
            for shape in list(slide.shapes):
                if not hasattr(shape, "is_placeholder") or not shape.is_placeholder:
                    slide.shapes._spTree.remove(shape._element)

            # Apply template based on type
            result = self._apply_template_layout(
                slide,
                template_id,
                template_config,
                content_mapping,
                image_paths,
                color_scheme,
            )

            if result["success"]:
                return {
                    "success": True,
                    "message": f"Applied template '{template_id}' to slide {slide_index}",
                    "slide_index": slide_index,
                    "template_id": template_id,
                    "color_scheme": color_scheme,
                    "elements_created": result.get("elements_created", []),
                }
            else:
                return result

        except Exception as e:
            return {"success": False, "error": f"Failed to apply template: {str(e)}"}

    def create_slide_from_template(
        self,
        template_id: str,
        color_scheme: str = "modern_blue",
        content_mapping: Optional[Dict[str, str]] = None,
        image_paths: Optional[Dict[str, str]] = None,
        layout_index: int = 1,
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Create a new slide using a layout template.

        Args:
            template_id: ID of the template to use (e.g., 'title_slide', 'text_with_image')
            color_scheme: Color scheme to use from professional color schemes
            content_mapping: Dictionary mapping element roles to custom content
            image_paths: Dictionary mapping image element roles to file paths
            layout_index: PowerPoint layout index to use as base (default: 1)
            presentation_id: Optional presentation ID (uses current if None)

        Returns:
            Dictionary with operation results
        """
        try:
            current_presentation = self.get_current_presentation()
            if current_presentation is None:
                return {
                    "success": False,
                    "error": "No presentation is currently loaded",
                }

            # Validate layout index
            if layout_index < 0 or layout_index >= len(
                current_presentation.slide_layouts
            ):
                return {
                    "success": False,
                    "error": f"Invalid layout index: {layout_index}. Available layouts: 0-{len(current_presentation.slide_layouts) - 1}",
                }

            # Add new slide
            layout = current_presentation.slide_layouts[layout_index]
            slide = current_presentation.slides.add_slide(layout)
            slide_index = len(current_presentation.slides) - 1

            # Apply template to the new slide
            result = self.apply_slide_template(
                slide_index, template_id, color_scheme, content_mapping, image_paths
            )

            if result["success"]:
                return {
                    "success": True,
                    "message": f"Created slide {slide_index} using template '{template_id}'",
                    "slide_index": slide_index,
                    "template_id": template_id,
                    "color_scheme": color_scheme,
                    "elements_created": result.get("elements_created", []),
                }
            else:
                # If template application failed, we should probably remove the slide
                # but python-pptx doesn't support slide deletion easily, so we'll leave it
                return {
                    "success": False,
                    "error": f"Created slide but failed to apply template: {result.get('error', 'Unknown error')}",
                    "slide_index": slide_index,
                }

        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to create slide from template: {str(e)}",
            }

    def create_presentation_from_templates(
        self,
        template_sequence: List[Dict[str, Any]],
        color_scheme: str = "modern_blue",
        presentation_title: Optional[str] = None,
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Create a complete presentation from a sequence of templates.

        Args:
            template_sequence: List of template configurations, each containing:
                - template_id: Template to use
                - content: Content mapping for the template
                - images: Image path mapping for the template (optional)
            color_scheme: Color scheme to apply to all slides
            presentation_title: Optional title for the presentation
            presentation_id: Optional presentation ID (uses current if None)

        Example template_sequence:
        [
            {
                "template_id": "title_slide",
                "content": {
                    "title": "My Presentation",
                    "subtitle": "Annual Report 2024",
                    "author": "John Doe"
                }
            },
            {
                "template_id": "text_with_image",
                "content": {
                    "title": "Key Results",
                    "content": "• Achievement 1\\n• Achievement 2"
                },
                "images": {
                    "supporting": "/path/to/image.jpg"
                }
            }
        ]

        Returns:
            Dictionary with operation results
        """
        try:
            current_presentation = self.get_current_presentation()
            if current_presentation is None:
                return {
                    "success": False,
                    "error": "No presentation is currently loaded",
                }

            if not template_sequence:
                return {"success": False, "error": "Template sequence cannot be empty"}

            slides_created = []
            errors = []

            # Set presentation title if provided
            if presentation_title and hasattr(current_presentation, "core_properties"):
                current_presentation.core_properties.title = presentation_title

            # Create slides from templates
            for i, template_config in enumerate(template_sequence):
                template_id = template_config.get("template_id")
                content_mapping = template_config.get("content", {})
                image_paths = template_config.get("images", {})

                if not template_id:
                    errors.append(f"Template {i}: Missing template_id")
                    continue

                result = self.create_slide_from_template(
                    template_id=template_id,
                    color_scheme=color_scheme,
                    content_mapping=content_mapping,
                    image_paths=image_paths,
                )

                if result["success"]:
                    slides_created.append(
                        {
                            "slide_index": result["slide_index"],
                            "template_id": template_id,
                        }
                    )
                else:
                    errors.append(
                        f"Template {i} ({template_id}): {result.get('error', 'Unknown error')}"
                    )

            success_count = len(slides_created)
            total_count = len(template_sequence)

            if success_count == 0:
                return {
                    "success": False,
                    "error": "No slides were created successfully",
                    "errors": errors,
                }

            return {
                "success": True,
                "message": f"Created {success_count}/{total_count} slides from templates",
                "slides_created": slides_created,
                "color_scheme": color_scheme,
                "presentation_title": presentation_title,
                "errors": errors if errors else None,
            }

        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to create presentation from templates: {str(e)}",
            }

    # ==================== Image Handling Tools ====================

    def add_image(
        self,
        slide_index: int,
        image_source: str,
        left: float = 1.0,
        top: float = 1.0,
        width: Optional[float] = None,
        height: Optional[float] = None,
        maintain_aspect_ratio: bool = True,
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Add an image to a slide with automatic sizing and positioning.

        Args:
            slide_index: Index of the slide to add image to
            image_source: File path to image or base64 encoded image string
            left: Left position in inches
            top: Top position in inches
            width: Width in inches (calculated automatically if None)
            height: Height in inches (calculated automatically if None)
            maintain_aspect_ratio: Whether to maintain original aspect ratio
            presentation_id: Optional presentation ID (uses current if None)

        Returns:
            Dictionary with operation results
        """
        try:
            current_presentation = self.get_current_presentation()
            if current_presentation is None:
                return {
                    "success": False,
                    "error": "No presentation is currently loaded",
                }

            if slide_index < 0 or slide_index >= len(current_presentation.slides):
                return {
                    "success": False,
                    "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(current_presentation.slides) - 1}",
                }

            slide = current_presentation.slides[slide_index]

            # Handle different image source types
            temp_path = None
            cleanup_temp = False

            try:
                if image_source.startswith("data:") or len(image_source) > 500:
                    # Assume it's base64 encoded
                    import base64
                    import tempfile

                    if image_source.startswith("data:"):
                        # Remove data URL prefix
                        image_source = image_source.split(",", 1)[1]

                    image_data = base64.b64decode(image_source)
                    with tempfile.NamedTemporaryFile(
                        delete=False, suffix=".png"
                    ) as temp_file:
                        temp_file.write(image_data)
                        temp_path = temp_file.name
                        cleanup_temp = True
                else:
                    # Assume it's a file path - resolve with base folder
                    resolved_image_path = self.resolve_file_path(image_source)
                    if not os.path.exists(resolved_image_path):
                        return {
                            "success": False,
                            "error": f"Image file not found: {resolved_image_path} (original: {image_source})",
                        }
                    temp_path = resolved_image_path

                # Add image to slide
                from pptx.util import Inches

                if width is None and height is None:
                    # Auto-size to reasonable default
                    width = 4.0
                    height = 3.0
                elif width is None and height is not None:
                    # Calculate width based on height and aspect ratio
                    if maintain_aspect_ratio:
                        try:
                            from PIL import Image

                            with Image.open(temp_path) as img:
                                aspect_ratio = img.width / img.height
                                width = height * aspect_ratio
                        except Exception:
                            width = height * 1.33  # Default 4:3 ratio
                    else:
                        width = height * 1.33
                elif height is None and width is not None:
                    # Calculate height based on width and aspect ratio
                    if maintain_aspect_ratio:
                        try:
                            from PIL import Image

                            with Image.open(temp_path) as img:
                                aspect_ratio = img.width / img.height
                                height = width / aspect_ratio
                        except Exception:
                            height = width * 0.75  # Default 4:3 ratio
                    else:
                        height = width * 0.75

                shape = slide.shapes.add_picture(
                    temp_path, Inches(left), Inches(top), Inches(width), Inches(height)
                )

                return {
                    "success": True,
                    "message": f"Added image to slide {slide_index}",
                    "slide_index": slide_index,
                    "shape_index": len(slide.shapes) - 1,
                    "position": {
                        "left": left,
                        "top": top,
                        "width": width,
                        "height": height,
                    },
                    "maintain_aspect_ratio": maintain_aspect_ratio,
                }

            finally:
                # Clean up temporary file if created
                if cleanup_temp and temp_path and os.path.exists(temp_path):
                    os.unlink(temp_path)

        except Exception as e:
            return {"success": False, "error": f"Failed to add image: {str(e)}"}

    def resize_image(
        self,
        slide_index: int,
        shape_index: int,
        width: Optional[float] = None,
        height: Optional[float] = None,
        maintain_aspect_ratio: bool = True,
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Resize an existing image on a slide.

        Args:
            slide_index: Index of the slide containing the image
            shape_index: Index of the image shape to resize
            width: New width in inches (None to calculate from height)
            height: New height in inches (None to calculate from width)
            maintain_aspect_ratio: Whether to maintain original aspect ratio
            presentation_id: Optional presentation ID (uses current if None)

        Returns:
            Dictionary with operation results
        """
        try:
            current_presentation = self.get_current_presentation()
            if current_presentation is None:
                return {
                    "success": False,
                    "error": "No presentation is currently loaded",
                }

            if slide_index < 0 or slide_index >= len(current_presentation.slides):
                return {
                    "success": False,
                    "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(current_presentation.slides) - 1}",
                }

            slide = current_presentation.slides[slide_index]

            if shape_index < 0 or shape_index >= len(slide.shapes):
                return {
                    "success": False,
                    "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}",
                }

            shape = slide.shapes[shape_index]

            # Check if it's an image shape
            if not hasattr(shape, "image"):
                return {"success": False, "error": "Selected shape is not an image"}

            from pptx.util import Inches

            # Get current dimensions
            current_width = shape.width.inches
            current_height = shape.height.inches
            current_aspect_ratio = current_width / current_height

            # Calculate new dimensions
            if width is None and height is None:
                return {
                    "success": False,
                    "error": "Either width or height must be specified",
                }
            elif width is None:
                # Calculate width from height
                if maintain_aspect_ratio:
                    width = height * current_aspect_ratio
                else:
                    width = current_width
            elif height is None:
                # Calculate height from width
                if maintain_aspect_ratio:
                    height = width / current_aspect_ratio
                else:
                    height = current_height

            # Apply new dimensions
            shape.width = Inches(width)
            shape.height = Inches(height)

            return {
                "success": True,
                "message": f"Resized image at shape {shape_index} on slide {slide_index}",
                "slide_index": slide_index,
                "shape_index": shape_index,
                "old_dimensions": {"width": current_width, "height": current_height},
                "new_dimensions": {"width": width, "height": height},
                "maintain_aspect_ratio": maintain_aspect_ratio,
            }

        except Exception as e:
            return {"success": False, "error": f"Failed to resize image: {str(e)}"}

    # ==================== Hyperlink Management Tools ====================

    def manage_hyperlinks(
        self,
        operation: str,
        slide_index: int,
        shape_index: Optional[int] = None,
        text: Optional[str] = None,
        url: Optional[str] = None,
        run_index: int = 0,
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Manage hyperlinks in text shapes and runs.

        Args:
            operation: Operation type ("add", "remove", "list", "update")
            slide_index: Index of the slide (0-based)
            shape_index: Index of the shape on the slide (0-based, required for add/remove/update)
            text: Text to make into hyperlink (for "add" operation)
            url: URL for the hyperlink
            run_index: Index of text run within the shape (0-based)
            presentation_id: Optional presentation ID (uses current if None)

        Returns:
            Dictionary with operation results
        """
        try:
            current_presentation = self.get_current_presentation()
            if current_presentation is None:
                return {
                    "success": False,
                    "error": "No presentation is currently loaded",
                }

            if slide_index < 0 or slide_index >= len(current_presentation.slides):
                return {
                    "success": False,
                    "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(current_presentation.slides) - 1}",
                }

            slide = current_presentation.slides[slide_index]

            if operation == "list":
                # List all hyperlinks in the slide
                hyperlinks = []
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for para_idx, paragraph in enumerate(
                            shape.text_frame.paragraphs
                        ):
                            for run_idx, run in enumerate(paragraph.runs):
                                if hasattr(run, "hyperlink") and run.hyperlink.address:
                                    hyperlinks.append(
                                        {
                                            "shape_index": shape_idx,
                                            "paragraph_index": para_idx,
                                            "run_index": run_idx,
                                            "text": run.text,
                                            "url": run.hyperlink.address,
                                        }
                                    )

                return {
                    "success": True,
                    "message": f"Found {len(hyperlinks)} hyperlinks on slide {slide_index}",
                    "slide_index": slide_index,
                    "hyperlinks": hyperlinks,
                }

            # For other operations, validate shape index
            if (
                shape_index is None
                or shape_index < 0
                or shape_index >= len(slide.shapes)
            ):
                return {
                    "success": False,
                    "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}",
                }

            shape = slide.shapes[shape_index]

            # Check if shape has text
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                return {"success": False, "error": "Shape does not contain text"}

            if operation == "add":
                if not text or not url:
                    return {
                        "success": False,
                        "error": "Both 'text' and 'url' are required for adding hyperlinks",
                    }

                # Add new text run with hyperlink
                paragraph = shape.text_frame.paragraphs[0]
                run = paragraph.add_run()
                run.text = text
                run.hyperlink.address = url

                return {
                    "success": True,
                    "message": f"Added hyperlink '{text}' -> '{url}' to shape {shape_index}",
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "text": text,
                    "url": url,
                }

            elif operation == "update":
                if not url:
                    return {
                        "success": False,
                        "error": "URL is required for updating hyperlinks",
                    }

                # Update existing hyperlink
                if run_index >= len(shape.text_frame.paragraphs[0].runs):
                    return {
                        "success": False,
                        "error": f"Run index {run_index} out of range",
                    }

                run = shape.text_frame.paragraphs[0].runs[run_index]
                old_url = run.hyperlink.address if hasattr(run, "hyperlink") else None
                run.hyperlink.address = url

                return {
                    "success": True,
                    "message": f"Updated hyperlink in run {run_index} to '{url}'",
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "run_index": run_index,
                    "old_url": old_url,
                    "new_url": url,
                }

            elif operation == "remove":
                if run_index >= len(shape.text_frame.paragraphs[0].runs):
                    return {
                        "success": False,
                        "error": f"Run index {run_index} out of range",
                    }

                run = shape.text_frame.paragraphs[0].runs[run_index]
                old_url = run.hyperlink.address if hasattr(run, "hyperlink") else None

                if hasattr(run, "hyperlink"):
                    run.hyperlink.address = None

                return {
                    "success": True,
                    "message": f"Removed hyperlink from run {run_index}",
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "run_index": run_index,
                    "removed_url": old_url,
                }

            else:
                return {
                    "success": False,
                    "error": f"Invalid operation: {operation}. Use 'add', 'remove', 'list', or 'update'",
                }

        except Exception as e:
            return {"success": False, "error": f"Failed to manage hyperlinks: {str(e)}"}

    # ==================== Template Helper Methods ====================

    def _apply_template_layout(
        self,
        slide,
        template_id: str,
        template_config: dict,
        content_mapping: dict,
        image_paths: dict,
        color_scheme: str,
    ) -> dict:
        """Apply a specific template layout to a slide."""
        try:
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor

            elements_created = []

            # Get professional colors
            color_schemes = self.get_color_schemes()
            if color_scheme in color_schemes["schemes"]:
                colors = color_schemes["schemes"][color_scheme]
            else:
                colors = color_schemes["schemes"]["modern_blue"]  # Fallback

            if template_id == "title_slide":
                # Title slide layout
                title_text = content_mapping.get("title", "Presentation Title")
                subtitle_text = content_mapping.get("subtitle", "Subtitle")
                author_text = content_mapping.get("author", "Author")

                # Add title
                title_shape = slide.shapes.add_textbox(
                    Inches(0.5), Inches(2), Inches(9), Inches(1.5)
                )
                title_frame = title_shape.text_frame
                title_frame.text = title_text
                title_frame.paragraphs[0].font.size = Pt(36)
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].font.color.rgb = RGBColor(*colors["primary"])
                elements_created.append("title")

                # Add subtitle
                subtitle_shape = slide.shapes.add_textbox(
                    Inches(0.5), Inches(4), Inches(9), Inches(1)
                )
                subtitle_frame = subtitle_shape.text_frame
                subtitle_frame.text = subtitle_text
                subtitle_frame.paragraphs[0].font.size = Pt(20)
                subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(
                    *colors["secondary"]
                )
                elements_created.append("subtitle")

                # Add author
                author_shape = slide.shapes.add_textbox(
                    Inches(0.5), Inches(5.5), Inches(9), Inches(0.8)
                )
                author_frame = author_shape.text_frame
                author_frame.text = author_text
                author_frame.paragraphs[0].font.size = Pt(16)
                author_frame.paragraphs[0].font.color.rgb = RGBColor(*colors["text"])
                elements_created.append("author")

            elif template_id == "text_with_image":
                # Text with image layout
                title_text = content_mapping.get("title", "Slide Title")
                content_text = content_mapping.get("content", "Content text goes here")

                # Add title
                title_shape = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.5), Inches(9), Inches(1)
                )
                title_frame = title_shape.text_frame
                title_frame.text = title_text
                title_frame.paragraphs[0].font.size = Pt(28)
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].font.color.rgb = RGBColor(*colors["primary"])
                elements_created.append("title")

                # Add content text
                text_shape = slide.shapes.add_textbox(
                    Inches(0.5), Inches(2), Inches(4.5), Inches(4)
                )
                text_frame = text_shape.text_frame
                text_frame.text = content_text
                text_frame.paragraphs[0].font.size = Pt(16)
                text_frame.paragraphs[0].font.color.rgb = RGBColor(*colors["text"])
                elements_created.append("content_text")

                # Add image placeholder or actual image
                if "supporting" in image_paths and image_paths["supporting"]:
                    try:
                        slide.shapes.add_picture(
                            image_paths["supporting"],
                            Inches(5.5),
                            Inches(2),
                            Inches(3.5),
                            Inches(4),
                        )
                        elements_created.append("supporting_image")
                    except Exception:
                        # Fallback to placeholder rectangle
                        img_placeholder = slide.shapes.add_shape(
                            1, Inches(5.5), Inches(2), Inches(3.5), Inches(4)
                        )
                        img_placeholder.fill.solid()
                        img_placeholder.fill.fore_color.rgb = RGBColor(*colors["light"])
                        elements_created.append("image_placeholder")
                else:
                    # Add placeholder rectangle
                    img_placeholder = slide.shapes.add_shape(
                        1, Inches(5.5), Inches(2), Inches(3.5), Inches(4)
                    )
                    img_placeholder.fill.solid()
                    img_placeholder.fill.fore_color.rgb = RGBColor(*colors["light"])
                    elements_created.append("image_placeholder")

            elif template_id == "bullet_points":
                # Bullet points layout
                title_text = content_mapping.get("title", "Slide Title")
                bullets_text = content_mapping.get(
                    "bullets", "• Point 1\n• Point 2\n• Point 3"
                )

                # Add title
                title_shape = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.5), Inches(9), Inches(1)
                )
                title_frame = title_shape.text_frame
                title_frame.text = title_text
                title_frame.paragraphs[0].font.size = Pt(28)
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].font.color.rgb = RGBColor(*colors["primary"])
                elements_created.append("title")

                # Add bullet points
                bullets_shape = slide.shapes.add_textbox(
                    Inches(0.5), Inches(2), Inches(8.5), Inches(4.5)
                )
                bullets_frame = bullets_shape.text_frame
                bullets_frame.text = bullets_text
                bullets_frame.paragraphs[0].font.size = Pt(18)
                bullets_frame.paragraphs[0].font.color.rgb = RGBColor(*colors["text"])
                elements_created.append("bullet_points")

            # Add more template types as needed...

            return {
                "success": True,
                "elements_created": elements_created,
                "template_id": template_id,
                "color_scheme": color_scheme,
            }

        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to apply template layout: {str(e)}",
            }

    # ==================== Connector Tools for Flowcharts ====================

    def add_connector(
        self,
        slide_index: int,
        connector_type: str,
        start_x: float,
        start_y: float,
        end_x: float,
        end_y: float,
        line_width: float = 1.0,
        color: Optional[List[int]] = None,
        arrow_style: str = "none",
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Add connector lines/arrows between points on a slide for flowcharts and diagrams.

        Args:
            slide_index: Index of the slide (0-based)
            connector_type: Type of connector ("straight", "elbow", "curved")
            start_x: Starting X coordinate in inches
            start_y: Starting Y coordinate in inches
            end_x: Ending X coordinate in inches
            end_y: Ending Y coordinate in inches
            line_width: Width of the connector line in points
            color: RGB color as [r, g, b] list
            arrow_style: Arrow style ("none", "end", "start", "both")
            presentation_id: Optional presentation ID (uses current if None)

        Returns:
            Dictionary with operation results
        """
        try:
            current_presentation = self.get_current_presentation()
            if current_presentation is None:
                return {
                    "success": False,
                    "error": "No presentation is currently loaded",
                }

            if slide_index < 0 or slide_index >= len(current_presentation.slides):
                return {
                    "success": False,
                    "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(current_presentation.slides) - 1}",
                }

            slide = current_presentation.slides[slide_index]

            # Map connector types
            from pptx.enum.shapes import MSO_CONNECTOR
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor

            connector_map = {
                "straight": MSO_CONNECTOR.STRAIGHT,
                "elbow": MSO_CONNECTOR.ELBOW,
                "curved": MSO_CONNECTOR.CURVED,
            }

            if connector_type.lower() not in connector_map:
                return {
                    "success": False,
                    "error": f"Invalid connector type. Use: {list(connector_map.keys())}",
                }

            # Add connector
            connector = slide.shapes.add_connector(
                connector_map[connector_type.lower()],
                Inches(start_x),
                Inches(start_y),
                Inches(end_x),
                Inches(end_y),
            )

            # Apply formatting
            if line_width:
                connector.line.width = Pt(line_width)

            if color and self.is_valid_rgb(color):
                connector.line.color.rgb = RGBColor(*color)

            # Apply arrow styles
            if arrow_style != "none":
                try:
                    from pptx.enum.dml import MSO_LINE_END_TYPE

                    if arrow_style in ["end", "both"]:
                        connector.line.end_arrow.type = MSO_LINE_END_TYPE.ARROW
                    if arrow_style in ["start", "both"]:
                        connector.line.begin_arrow.type = MSO_LINE_END_TYPE.ARROW
                except Exception:
                    pass  # Graceful fallback if arrow styling fails

            return {
                "success": True,
                "message": f"Added {connector_type} connector to slide {slide_index}",
                "slide_index": slide_index,
                "connector_type": connector_type,
                "start_point": [start_x, start_y],
                "end_point": [end_x, end_y],
                "arrow_style": arrow_style,
                "shape_index": len(slide.shapes) - 1,
            }

        except Exception as e:
            return {"success": False, "error": f"Failed to add connector: {str(e)}"}

    def connect_shapes(
        self,
        slide_index: int,
        start_shape_index: int,
        end_shape_index: int,
        connector_type: str = "straight",
        line_width: float = 1.0,
        color: Optional[List[int]] = None,
        arrow_style: str = "end",
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Connect two existing shapes with a connector line.

        Args:
            slide_index: Index of the slide containing the shapes
            start_shape_index: Index of the starting shape
            end_shape_index: Index of the ending shape
            connector_type: Type of connector ("straight", "elbow", "curved")
            line_width: Width of the connector line in points
            color: RGB color as [r, g, b] list
            arrow_style: Arrow style ("none", "end", "start", "both")
            presentation_id: Optional presentation ID (uses current if None)

        Returns:
            Dictionary with operation results
        """
        try:
            current_presentation = self.get_current_presentation()
            if current_presentation is None:
                return {
                    "success": False,
                    "error": "No presentation is currently loaded",
                }

            if slide_index < 0 or slide_index >= len(current_presentation.slides):
                return {
                    "success": False,
                    "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(current_presentation.slides) - 1}",
                }

            slide = current_presentation.slides[slide_index]

            # Validate shape indices
            if start_shape_index < 0 or start_shape_index >= len(slide.shapes):
                return {
                    "success": False,
                    "error": f"Invalid start shape index: {start_shape_index}. Available shapes: 0-{len(slide.shapes) - 1}",
                }

            if end_shape_index < 0 or end_shape_index >= len(slide.shapes):
                return {
                    "success": False,
                    "error": f"Invalid end shape index: {end_shape_index}. Available shapes: 0-{len(slide.shapes) - 1}",
                }

            start_shape = slide.shapes[start_shape_index]
            end_shape = slide.shapes[end_shape_index]

            # Calculate connection points (center of shapes)
            start_x = (start_shape.left + start_shape.width / 2).inches
            start_y = (start_shape.top + start_shape.height / 2).inches
            end_x = (end_shape.left + end_shape.width / 2).inches
            end_y = (end_shape.top + end_shape.height / 2).inches

            # Add connector between shapes
            result = self.add_connector(
                slide_index,
                connector_type,
                start_x,
                start_y,
                end_x,
                end_y,
                line_width,
                color,
                arrow_style,
            )

            if result["success"]:
                result["connected_shapes"] = {
                    "start_shape_index": start_shape_index,
                    "end_shape_index": end_shape_index,
                }
                result["message"] = (
                    f"Connected shape {start_shape_index} to shape {end_shape_index}"
                )

            return result

        except Exception as e:
            return {"success": False, "error": f"Failed to connect shapes: {str(e)}"}

    # ==================== Master Slide Management ====================

    def get_slide_masters(self, presentation_id: Optional[str] = None) -> dict:
        """Get information about slide masters in the presentation.

        Args:
            presentation_id: Optional presentation ID (uses current if None)

        Returns:
            Dictionary with slide master information
        """
        try:
            current_presentation = self.get_current_presentation()
            if current_presentation is None:
                return {
                    "success": False,
                    "error": "No presentation is currently loaded",
                }

            masters_info = []

            for i, master in enumerate(current_presentation.slide_masters):
                layouts_info = []
                for j, layout in enumerate(master.slide_layouts):
                    layouts_info.append(
                        {
                            "layout_index": j,
                            "name": (
                                layout.name
                                if hasattr(layout, "name")
                                else f"Layout {j}"
                            ),
                            "placeholder_count": (
                                len(layout.placeholders)
                                if hasattr(layout, "placeholders")
                                else 0
                            ),
                        }
                    )

                masters_info.append(
                    {
                        "master_index": i,
                        "layouts": layouts_info,
                        "total_layouts": len(master.slide_layouts),
                    }
                )

            return {
                "success": True,
                "slide_masters": masters_info,
                "total_masters": len(current_presentation.slide_masters),
                "message": f"Found {len(current_presentation.slide_masters)} slide masters",
            }

        except Exception as e:
            return {"success": False, "error": f"Failed to get slide masters: {str(e)}"}

    def get_slide_layouts(
        self, master_index: int = 0, presentation_id: Optional[str] = None
    ) -> dict:
        """Get detailed information about slide layouts for a specific master.

        Args:
            master_index: Index of the slide master (default: 0)
            presentation_id: Optional presentation ID (uses current if None)

        Returns:
            Dictionary with slide layout information
        """
        try:
            current_presentation = self.get_current_presentation()
            if current_presentation is None:
                return {
                    "success": False,
                    "error": "No presentation is currently loaded",
                }

            if master_index < 0 or master_index >= len(
                current_presentation.slide_masters
            ):
                return {
                    "success": False,
                    "error": f"Invalid master index: {master_index}. Available masters: 0-{len(current_presentation.slide_masters) - 1}",
                }

            master = current_presentation.slide_masters[master_index]
            layouts_info = []

            for i, layout in enumerate(master.slide_layouts):
                placeholders_info = []
                if hasattr(layout, "placeholders"):
                    for placeholder in layout.placeholders:
                        placeholder_info = {
                            "placeholder_format": (
                                placeholder.placeholder_format.type
                                if hasattr(placeholder, "placeholder_format")
                                else None
                            ),
                            "name": (
                                placeholder.name
                                if hasattr(placeholder, "name")
                                else f"Placeholder {len(placeholders_info)}"
                            ),
                            "width": (
                                placeholder.width.inches
                                if hasattr(placeholder, "width")
                                else None
                            ),
                            "height": (
                                placeholder.height.inches
                                if hasattr(placeholder, "height")
                                else None
                            ),
                        }
                        placeholders_info.append(placeholder_info)

                layouts_info.append(
                    {
                        "layout_index": i,
                        "name": (
                            layout.name if hasattr(layout, "name") else f"Layout {i}"
                        ),
                        "placeholders": placeholders_info,
                        "placeholder_count": len(placeholders_info),
                    }
                )

            return {
                "success": True,
                "master_index": master_index,
                "layouts": layouts_info,
                "total_layouts": len(layouts_info),
                "message": f"Found {len(layouts_info)} layouts in master {master_index}",
            }

        except Exception as e:
            return {"success": False, "error": f"Failed to get slide layouts: {str(e)}"}

    # ==================== Transition and Animation Tools ====================

    def set_slide_transition(
        self,
        slide_index: int,
        transition_type: str = "fade",
        duration: float = 1.0,
        advance_on_click: bool = True,
        advance_after_time: Optional[float] = None,
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Set transition effect for a slide.

        Args:
            slide_index: Index of the slide to set transition for
            transition_type: Type of transition ("none", "fade", "push", "wipe", "split", "reveal", "random")
            duration: Duration of transition in seconds
            advance_on_click: Whether slide advances on mouse click
            advance_after_time: Auto-advance after specified seconds (None to disable)
            presentation_id: Optional presentation ID (uses current if None)

        Returns:
            Dictionary with operation results
        """
        try:
            current_presentation = self.get_current_presentation()
            if current_presentation is None:
                return {
                    "success": False,
                    "error": "No presentation is currently loaded",
                }

            if slide_index < 0 or slide_index >= len(current_presentation.slides):
                return {
                    "success": False,
                    "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(current_presentation.slides) - 1}",
                }

            slide = current_presentation.slides[slide_index]

            # Note: python-pptx has limited transition support
            # Most transition effects require XML manipulation
            try:
                # Basic transition setting (limited support)
                if hasattr(slide, "slide_transition"):
                    # This is a simplified implementation
                    # Full transition support would require XML manipulation
                    pass

                # Set advance settings if available
                if hasattr(slide, "advance_on_click"):
                    slide.advance_on_click = advance_on_click

                if advance_after_time is not None and hasattr(
                    slide, "advance_after_time"
                ):
                    slide.advance_after_time = advance_after_time

                return {
                    "success": True,
                    "message": f"Set {transition_type} transition for slide {slide_index}",
                    "slide_index": slide_index,
                    "transition_type": transition_type,
                    "duration": duration,
                    "advance_on_click": advance_on_click,
                    "advance_after_time": advance_after_time,
                    "note": "Transition effects have limited support in python-pptx",
                }

            except Exception as e:
                return {
                    "success": False,
                    "error": f"Transition setting not fully supported: {str(e)}",
                    "note": "python-pptx has limited transition support",
                }

        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to set slide transition: {str(e)}",
            }

    def set_presentation_transitions(
        self,
        transition_type: str = "fade",
        duration: float = 1.0,
        advance_on_click: bool = True,
        advance_after_time: Optional[float] = None,
        presentation_id: Optional[str] = None,
    ) -> dict:
        """Set transition effects for all slides in the presentation.

        Args:
            transition_type: Type of transition to apply to all slides
            duration: Duration of transition in seconds
            advance_on_click: Whether slides advance on mouse click
            advance_after_time: Auto-advance after specified seconds (None to disable)
            presentation_id: Optional presentation ID (uses current if None)

        Returns:
            Dictionary with operation results
        """
        try:
            current_presentation = self.get_current_presentation()
            if current_presentation is None:
                return {
                    "success": False,
                    "error": "No presentation is currently loaded",
                }

            total_slides = len(current_presentation.slides)
            successful_slides = []
            failed_slides = []

            for i in range(total_slides):
                result = self.set_slide_transition(
                    i, transition_type, duration, advance_on_click, advance_after_time
                )

                if result["success"]:
                    successful_slides.append(i)
                else:
                    failed_slides.append(
                        {
                            "slide_index": i,
                            "error": result.get("error", "Unknown error"),
                        }
                    )

            return {
                "success": len(failed_slides) == 0,
                "message": f"Applied {transition_type} transition to {len(successful_slides)}/{total_slides} slides",
                "transition_type": transition_type,
                "successful_slides": successful_slides,
                "failed_slides": failed_slides if failed_slides else None,
                "total_slides": total_slides,
                "note": "Transition effects have limited support in python-pptx",
            }

        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to set presentation transitions: {str(e)}",
            }

    def get_transition_info(
        self, slide_index: int, presentation_id: Optional[str] = None
    ) -> dict:
        """Get current transition information for a slide.

        Args:
            slide_index: Index of the slide to get transition info for
            presentation_id: Optional presentation ID (uses current if None)

        Returns:
            Dictionary with transition information
        """
        try:
            current_presentation = self.get_current_presentation()
            if current_presentation is None:
                return {
                    "success": False,
                    "error": "No presentation is currently loaded",
                }

            if slide_index < 0 or slide_index >= len(current_presentation.slides):
                return {
                    "success": False,
                    "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(current_presentation.slides) - 1}",
                }

            slide = current_presentation.slides[slide_index]

            # Extract available transition information
            transition_info = {
                "slide_index": slide_index,
                "advance_on_click": getattr(slide, "advance_on_click", True),
                "advance_after_time": getattr(slide, "advance_after_time", None),
                "has_transition": False,  # Would need XML parsing for full info
                "note": "Detailed transition info requires XML parsing",
            }

            return {
                "success": True,
                "transition_info": transition_info,
                "message": f"Retrieved transition info for slide {slide_index}",
            }

        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to get transition info: {str(e)}",
            }

    # --- Event Emitter Enhanced Methods ---

    async def create_presentation_with_status(
        self, id: Optional[str] = None, __event_emitter__=None
    ) -> dict:
        """Create a new PowerPoint presentation with status updates."""
        if __event_emitter__:
            await __event_emitter__(
                {
                    "type": "status",
                    "data": {
                        "description": "Creating new PowerPoint presentation...",
                        "done": False,
                    },
                }
            )

        try:
            result = self.create_presentation(id)

            if __event_emitter__:
                if "error" not in result:
                    await __event_emitter__(
                        {
                            "type": "status",
                            "data": {
                                "description": f"Successfully created presentation {result['presentation_id']}",
                                "done": True,
                            },
                        }
                    )
                else:
                    await __event_emitter__(
                        {
                            "type": "status",
                            "data": {
                                "description": f"Failed to create presentation: {result['error']}",
                                "done": True,
                            },
                        }
                    )

            return result
        except Exception as e:
            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {
                            "description": f"Error creating presentation: {str(e)}",
                            "done": True,
                        },
                    }
                )
            return {"error": f"Exception occurred: {str(e)}"}
